# In-place editor for the revA/revB decks (2026-07-12, round 2). The user
# added comments directly to the built revA/revB files, so those files are
# now the bases - rebuilding from the main deck would destroy the comments.
# This script edits them in place:
#   both decks: comment 125 response (18-state deployment dotplot) and the
#               strata 2024 numbers/figure once the year-swap run lands
#   revB only:  condense wordy slides to one-point bullets (originals to
#               notes) and reorder for logical flow per the slide-2 outline
#
#   python methods/apply_deck_revisions_inplace.py
import copy
import csv
import os
from pptx import Presentation
from pptx.util import Inches, Pt

REV_A = "lessons_getting_more_signal_from_snap_qc_data_revA.pptx"
REV_B = "lessons_getting_more_signal_from_snap_qc_data_revB.pptx"
DEPLOY = "methods/state_similarity_v2/transfer_benchmark_train2223_test24"
STRATA = "methods/compare_hh_strata_v2/yearswap_train2223_test24"


def find_slide(p, substr):
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return s
    raise KeyError(f"no slide containing {substr!r}")


def find_tf(s, substr):
    for sh in s.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            return sh.text_frame
    raise KeyError(f"no shape containing {substr!r}")


def set_para(tf, match, new_text):
    for para in tf.paragraphs:
        if match in para.text:
            if para.runs:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.text = new_text
            return para
    raise KeyError(f"no paragraph containing {match!r}")


def add_para_after(tf, match, new_text):
    from pptx.text.text import _Paragraph
    for para in tf.paragraphs:
        if match in para.text:
            new_p = copy.deepcopy(para._p)
            para._p.addnext(new_p)
            np = _Paragraph(new_p, para._parent)
            if np.runs:
                np.runs[0].text = new_text
                for r in np.runs[1:]:
                    r.text = ""
            else:
                np.text = new_text
            return np
    raise KeyError(f"no paragraph containing {match!r}")


def swap_picture(s, new_path):
    pics = [sh for sh in s.shapes if sh.shape_type == 13]
    pic = pics[0]
    box = (pic.left, pic.top, pic.width, pic.height)
    pic._element.getparent().remove(pic._element)
    s.shapes.add_picture(new_path, box[0], box[1], width=box[2], height=box[3])


def add_note(s, text):
    tf = s.notes_slide.notes_text_frame
    existing = tf.text
    tf.text = (existing + "\n\n" if existing.strip() else "") + \
        "[2026-07-12 rev2] " + text


def set_body_lines(tf, lines):
    from pptx.text.text import _Paragraph
    first = tf.paragraphs[0]
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    anchor = first._p
    for _ in lines[1:]:
        new_p = copy.deepcopy(first._p)
        anchor.addnext(new_p)
        anchor = new_p
    for para, line in zip(tf.paragraphs, lines):
        pp = _Paragraph(para._p, para._parent)
        if pp.runs:
            pp.runs[0].text = line
            for r in pp.runs[1:]:
                r.text = ""
        else:
            pp.text = line


def condense(p, title_anchor, body_anchor, lines, why):
    sl = find_slide(p, title_anchor)
    tf = find_tf(sl, body_anchor)
    original = tf.text
    set_body_lines(tf, lines)
    add_note(sl, "revB condensed (" + why + "). Original text: " +
             original.replace("\n", " | "))


# ── edits for BOTH decks ──────────────────────────────────────────────────────
def apply_common(p):
    # comment 125: "did we have an 18 state version of this?" (guarded)
    s22 = find_slide(p, "The national rules, deployed a year ahead")
    tf22 = find_tf(s22, "Every state clears its base rate")
    already = ("0.18-0.43" in tf22.text or
               "Mined on all states' 2022-23; applied in confidence order"
               in tf22.text)
    if not already:
        swap_picture(s22, os.path.join(DEPLOY, "workshop_national_dotplot_budget10.png"))
        set_para(tf22, "Every state clears its base rate",
                 "- Every state clears its base rate: precision 0.18-0.43 against "
                 "8-22% base rates, a 1.5-3.4x lift over random review.")
    if not already:
        add_note(s22, "Comment 'did we have an 18 state version? can we include "
                 "that?': yes - chart swapped to the 18-state workshop version "
                 "(workshop_national_dotplot_budget10.png; adds ME, MD, MO, MA, "
                 "DC, TN). Range updated: 10%-budget precision 0.177 (NJ) to "
                 "0.433 (DC); base rates 8.5-21.9%; lift 1.5x (TX) to 3.4x (WA).")

    # strata year-swap results (skipped gracefully until the run lands)
    st_path = os.path.join(STRATA, "strata_summary.csv")
    if os.path.exists(st_path):
        st = {r["scheme"]: r for r in csv.DictReader(open(st_path, newline=""))}
        pooled, coarse = st.get("Pooled (no split)"), st.get("1 / 2-3 / 4+")
        s14 = find_slide(p, "Split by coarse household-size strata")
        if pooled and coarse and "Re-tested on 2024" not in \
                find_tf(s14, "coverage parity").text:
            add_para_after(
                find_tf(s14, "coverage parity"), "coverage parity",
                "Re-tested on 2024: the split wins precision at matched recall "
                f"({coarse['mean_precision']} vs {pooled['mean_precision']} pooled); "
                f"pooled reaches further at the 0.20 floor "
                f"({pooled['recall_at_020']} vs {coarse['recall_at_020']}). "
                "The coarse split stays the safe default; 5-way ties it at "
                "~1.6x the compute.")
            fig = os.path.join(STRATA, "strata_sweeps.png")
            if os.path.exists(fig):
                swap_picture(s14, fig)
            add_note(s14, "2024 strata re-run landed - PARTIAL replication: "
                     "on 2023 pooling matched the split's precision (0.226 vs "
                     "0.222) and the split won reach; on 2024 the split wins "
                     "precision (0.302 vs 0.262) and pooling wins floor-reach "
                     "(0.844 vs 0.794). Consistent across years: the coarse "
                     "split never loses and stays the default. NOT replicated: "
                     "'5-way is worse' - on 2024 it ties the 3-way (0.304 vs "
                     "0.302) at ~1.6x compute, so the claim softens to 'no "
                     "better, costlier'. Figure swapped to the 2024 version.")


# ── revB only: condensation + reorder ─────────────────────────────────────────
def apply_b(p):
    condense(p, "Two data-build choices", "Cases whose review found", [
        "- 31% of error cases were being dropped (second error element); restoring them tripled the qualifying rules.",
        "- Rule content survived the fix: 93% of the old rule set persisted.",
        "- Blank deduction fields dropped ~16% of Washington's caseload; zero-fill + flag kept them.",
        "- Both defects were invisible in model metrics - found only by reconciling counts against the raw files."],
        "one point per bullet")
    condense(p, "Evaluate at review budgets", "Confidence floors produce", [
        "- A confidence floor produces whatever workload it produces - the 0.20 floor flags ~half the caseload.",
        "- So we report budget-filled performance: add rules in confidence order until 5% or 10% of caseload is flagged.",
        "- Deployment medians: 0.30 precision / 16% of error $ at 5%; 0.27 / 25% at 10% - vs 8-17% base rates."],
        "Mississippi floor-artifact example moved to notes")
    condense(p, "The national rules, deployed a year ahead", "The deployment test:", [
        "- Mined on all states' 2022-23; applied in confidence order to 10% of caseload; scored on 2024.",
        "- Every state clears its base rate: 1.5-3.4x lift over random review."],
        "include-own-state finding moved to notes")
    condense(p, "Donor-state similarity", "We compute state similarity", [
        "- Similarity computed from the QC microdata (which rules fire in a state; rare rules weighted up); 5 nearest donors.",
        "- Neighbor lists change across eras: only ~2 of 5 members persist.",
        "- Same-era, neighbor pools looked competitive at a 10% budget; tested on 2024 the advantage vanished.",
        "- Own-state-only mining is the high-variance option: best results anywhere (CT, VA), worst too (WA, below base rate)."],
        "full numbers moved to notes")
    condense(p, "Adapting the national rules to a state", "filtered - keep the national rules", [
        "- filtered: national rules re-qualified on the state's own data.",
        "- tuned: thresholds re-searched on the state's data; most defensible variant deployed.",
        "- hybrid: the settled grid-search scheme - careful gate, dollar-max pick.",
        "- All arms budget-filled on the state's 2024 at identical review volume."],
        "gate parameters moved to notes")
    condense(p, "adaptation does not beat the national order", "No adaptation scheme beats", [
        "- No adaptation scheme beats the national ordering for the median state.",
        "- Adaptation pays where the national list is weak (NJ, MS, DC); hurts where it is strong (MI, WA, MA).",
        "- The tuned arm is too conservative to deploy."],
        "win counts and per-state numbers moved to notes")
    condense(p, "5% review budget: national as-is vs adapted", "precision @ recall", [
        "- precision @ recall ($ share) on 2024, trained 2022-23. 'Rules qualified' counts rules, not cases."],
        "qualification + top-500 detail moved to notes")
    condense(p, "10% review budget: national as-is vs adapted", "precision @ recall", [
        "- precision @ recall ($ share) on 2024, trained 2022-23. 'Rules qualified' counts rules, not cases."],
        "qualification + top-500 detail moved to notes")
    condense(p, "The deployed list is a few dozen rules", "The budget fill scans", [
        "- The budget fill scans ~45k rules, but the union is BUILT by only 16-39 (5%) / 34-67 (10%) of them.",
        "- Redundant rules 'admit' at zero added cases - that is why admitted counts look huge.",
        "- The short list is reproducible from covariates alone; no outcomes needed."],
        "one point per bullet")
    condense(p, "The most widely deployed rules", "The three rules deployed", [
        "- The three most widely deployed rules per HH stratum (10% budget). Benefits-near-max and deduction levels dominate."],
        "exact feature counts moved to notes")
    condense(p, "The deliverable: one ranked rule list", "Built in advance from public data", [
        "- One ranked list per state: fill to the budget on the state's own caseload (core), keep filling to 3x (buffer).",
        "- The state activates rules in order while capacity fits - no outcome data at any step.",
        "- Workload lands on budget as patterns drift (raw core drifted 2.3-12%; walked lists land on target in all 18).",
        "- Median activated: 23 rules at 5%, 42 at 10%."],
        "buffer rationale + validation recipe moved to notes")
    condense(p, "Blending state and national rules", "Merge the state's own mined pool", [
        "- Rank EVERY rule - state or national - by its own 99% lower bound; the bound prices the certainty gap.",
        "- As a single recipe the blend beats national-only at 5% (0.324 vs 0.294) and ties at 10%.",
        "- Where state rules clear the bar, interleaving beats both parents (AZ, DC, MS).",
        "- Blind spot: national bounds say nothing about transfer - NJ's own rules never enter; the own-pool fallback exists."],
        "full mechanics moved to notes")

    # reorder to follow the slide-2 outline + learning order
    anchors = [
        "- modeling lessons",
        "Approach is to extract rules",
        "Intuition for new modeling approach",
        "What we ran",
        "Two data-build choices",
        "Limitations of the public QC data",
        "Selecting rules on raw training precision",
        "Better way to select rules",
        "Bigger ensembles widen the menu",
        "By error category vs all errors",
        "Tuning: what moved held-out performance",
        "More strata seemed to be a big win",
        "Split by coarse household-size strata",
        "Re-test your selection choices",
        "@DIVIDER_BARE",
        "state-generated and nationally-generated",
        "At state scale, confidence bounds alone",
        "For thin states, same-era neighbor data",
        "Donor-state similarity",
        "Evaluate at review budgets",
        "The national rules, deployed a year ahead",
        "Two-regime",
        "Adapting the national rules to a state",
        "adaptation does not beat the national order",
        "The adaptation schemes, drawn",
        "5% review budget: national as-is vs adapted",
        "10% review budget: national as-is vs adapted",
        "The deployed list is a few dozen rules",
        "The most widely deployed rules",
        "The deliverable: one ranked rule list",
        "Freeze each state's list in advance",
        "Blending state and national rules",
        "Blend vs national vs own rules, 10% budget",
        "What a state is handed",
        "live on github",
        "Appendix: one-at-a-time xgboost sweep",
    ]
    bare = None
    for sl in p.slides:
        joined = " ".join(sh.text_frame.text for sh in sl.shapes
                          if sh.has_text_frame)
        if "Optimizing rule sets for states" in joined and \
                "state-generated" not in joined:
            bare = sl
            break
    order = []
    for a in anchors:
        sl = bare if a == "@DIVIDER_BARE" else find_slide(p, a)
        order.append(sl)
    assert len(set(id(x) for x in order)) == len(order)
    id_map = {id(sl): sldid for sldid, sl in zip(p.slides._sldIdLst, p.slides)}
    ordered_ids = [id_map[id(sl)] for sl in order]
    remaining = [x for x in p.slides._sldIdLst if x not in ordered_ids]
    lst = p.slides._sldIdLst
    for el in list(lst):
        lst.remove(el)
    for el in ordered_ids + remaining:
        lst.append(el)


def run(path, is_b):
    p = Presentation(path)
    apply_common(p)
    if is_b:
        try:
            tf = find_tf(find_slide(p, "Two data-build choices"),
                         "31% of error cases were being dropped")
            b_done = True
        except KeyError:
            b_done = False
        if not b_done:
            apply_b(p)
    p.save(path)
    print(f"updated {path} ({len(p.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    run(REV_A, is_b=False)
    run(REV_B, is_b=True)
