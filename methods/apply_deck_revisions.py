# Builds two revision variants of the user-edited lessons deck (33 slides,
# re-read 2026-07-12 late edit):
#   revA - direct responses to the 11 margin comments, the alphabetical-chart
#          convention, and corrections of stale/misstated numbers
#   revB - revA plus suggested extras (kept minimal; short bullets only)
# Conventions honored: user's wording and deletions are never undone; slide
# text stays terse (detail goes to speaker notes, prefixed [2026-07-12]);
# slides are located by their text, not position, so further edits survive.
#
#   python methods/apply_deck_revisions.py
import copy
import csv
import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

SRC = "slides/lessons_getting_more_signal_from_snap_qc_data.pptx"
OUT_A = "slides/lessons_getting_more_signal_from_snap_qc_data_revA.pptx"
OUT_B = "slides/lessons_getting_more_signal_from_snap_qc_data_revB.pptx"
DEPLOY = "methods/state_similarity_v2/transfer_benchmark_train2223_test24"
STRATA = "methods/compare_hh_strata_v2/yearswap_train2223_test24"
FIG_2024_FLOORS = "presentation_figures/floor_definitions_educational_2024.png"


# ── generic helpers ───────────────────────────────────────────────────────────
def find_slide(p, substr):
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return s
    raise KeyError(f"no slide containing {substr!r}")


def find_tf(s, substr):
    for sh in s.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            return sh.text_frame
    raise KeyError(f"no shape containing {substr!r}")


def set_para(tf, match, new_text):
    for para in tf.paragraphs:
        if match in para.text:
            if para.runs:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.text = new_text
            return para
    raise KeyError(f"no paragraph containing {match!r}")


def add_para_after(tf, match, new_text):
    from pptx.text.text import _Paragraph
    for para in tf.paragraphs:
        if match in para.text:
            new_p = copy.deepcopy(para._p)
            para._p.addnext(new_p)
            np = _Paragraph(new_p, para._parent)
            if np.runs:
                np.runs[0].text = new_text
                for r in np.runs[1:]:
                    r.text = ""
            else:
                np.text = new_text
            return np
    raise KeyError(f"no paragraph containing {match!r}")


def swap_picture(s, new_path, which=0):
    pics = [sh for sh in s.shapes if sh.shape_type == 13]
    pic = pics[which]
    box = (pic.left, pic.top, pic.width, pic.height)
    pic._element.getparent().remove(pic._element)
    s.shapes.add_picture(new_path, box[0], box[1], width=box[2], height=box[3])


def add_note(s, text):
    tf = s.notes_slide.notes_text_frame
    existing = tf.text
    tf.text = (existing + "\n\n" if existing.strip() else "") + \
        "[2026-07-12 revision] " + text


def copy_shape(src_slide, dst_slide, substr):
    for sh in src_slide.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            el = copy.deepcopy(sh._element)
            dst_slide.shapes._spTree.append(el)
            for sh2 in dst_slide.shapes:
                if sh2._element is el:
                    return sh2
    raise KeyError(substr)


def move_slide(p, from_pos, to_pos):
    lst = p.slides._sldIdLst
    el = list(lst)[from_pos - 1]
    lst.remove(el)
    lst.insert(to_pos - 1, el)


def slide_pos(p, s):
    for i, x in enumerate(p.slides, 1):
        if x == s:
            return i
    raise KeyError


def blended_numbers():
    path = os.path.join(DEPLOY, "blended_frozen_results.csv")
    return [r for r in csv.DictReader(open(path, newline=""))
            if r["variant"] == "lcb99"]


def strata_numbers():
    path = os.path.join(STRATA, "strata_summary.csv")
    if not os.path.exists(path):
        return None
    return {r["scheme"]: r for r in csv.DictReader(open(path, newline=""))}


# ── variant A: the 11 comment responses + conventions + stale fixes ──────────
def apply_variant_a(p):
    # C135 (slide "Intuition for new modeling approach"): constrained-RF figure
    s3 = find_slide(p, "Intuition for new modeling approach")
    s3.shapes.add_picture("methods/parameter_tuning_v2/mtry_frontier.png",
                          Inches(6.0), Inches(0.85), width=Inches(3.7))
    add_note(s3, "Comment 'Could be uniform random forests / add figure for "
             "constrained RF': figure added (mtry frontier - 2 beats fully "
             "random 1 and looser 4; methods/visualize_mtry_frontier_v2.R). On "
             "uniform random forests: plausible alternative for the diversity "
             "role, but untested here - ranger with mtry=2 is what we measured; "
             "worth a head-to-head before claiming either way.")

    # C118 (untitled draft slide with the two typed-vs-pooled charts): finish it
    s8 = find_slide(p, "Selecting rules on raw training precision")
    s7 = None
    for s in p.slides:
        pics = [sh for sh in s.shapes if sh.shape_type == 13]
        texts = [sh for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if len(pics) == 2 and not texts:
            s7 = s
            break
    if s7 is not None:
        title = copy_shape(s8, s7, "Selecting rules on raw training precision")
        title.left, title.top = Inches(0.44), Inches(0.30)
        title.width, title.height = Inches(3.9), Inches(1.15)
        title.text_frame.word_wrap = True
        set_para(title.text_frame, "Selecting rules",
                 "By error category vs all errors: mine both, pool them")
        for para in title.text_frame.paragraphs:
            for r in para.runs:
                r.font.size = Pt(20)
        body = copy_shape(s8, s7, "V1 'train precision")
        body.left, body.top = Inches(0.44), Inches(1.60)
        body.width, body.height = Inches(3.9), Inches(1.7)
        tfb = body.text_frame
        tfb.word_wrap = True
        set_para(tfb, "V1 'train precision",
                 "Typed frames edge one all-errors model slightly; the UNION "
                 "wins: +3-6pp recall for ~1pp precision.")
        set_para(tfb, "But, rules selected on the HOLDOUT",
                 "Top: 2024 replication. Bottom: original (2023).")
        set_para(tfb, "And, before any selection", "")
        set_para(tfb, "So the decay is selection noise", "")
        for para in tfb.paragraphs:
            for r in para.runs:
                r.font.size = Pt(13)
        pics = sorted([sh for sh in s7.shapes if sh.shape_type == 13],
                      key=lambda sh: sh.top)
        for pic, path in zip(pics, [
                "methods/compare_anyerror_vs_typed_v2/yearswap_train2223_test24/anyerror_vs_typed_sweep.png",
                "methods/compare_anyerror_vs_typed_v2/anyerror_vs_typed_sweep.png"]):
            box = (pic.left, pic.top, pic.width, pic.height)
            pic._element.getparent().remove(pic._element)
            s7.shapes.add_picture(path, box[0], box[1],
                                  width=box[2], height=box[3])
        add_note(s7, "Comment 'add a slide showing by-category vs all-errors vs "
                 "pooled': completed this draft slide - title + two short "
                 "bullets; your two pasted charts kept (top = year-swap "
                 "replication, bottom = original). Numbers: findings section 3 "
                 "- combined beats typed alone by +3-6pp recall at fixed "
                 "floors, ~0.7-2pp precision cost, replicated on 2024.")

    # C11A (visibility slide): verified
    s6 = find_slide(p, "Limitations of the public QC data")
    add_note(s6, "Comment 'Let's check these numbers': verified against "
             "methods/state_error_accounting/visibility_by_state_2022_2024.csv "
             "- NJ 42.7%, TN 50.8%, AR/MO/UT 53.4%, WA 78.5%, VA 80.2%, LA "
             "80.5%; national (error-weighted) 71.1%; range 42.7 (NJ) to 90.6 "
             "(GA). All table cells match.")

    # C11C (LCB slide): examples + min-n + 2024 chart
    s9 = find_slide(p, "Better way to select rules")
    tf9 = find_tf(s9, "At 10 cases flagged")
    set_para(tf9, "At 10 cases flagged",
             "At 10 cases flagged, aiming for 0.3 precision: need 7 of 10 hits (raw 0.70)")
    set_para(tf9, "At 100 cases flagged",
             "At 100 cases flagged, aiming for 0.3 precision: need 41 of 100 hits (raw 0.41)")
    add_para_after(tf9, "At 100 cases flagged",
                   "(Raw selection also had a support floor - n >= 10 - and the "
                   "curse survived it.)")
    if os.path.exists(FIG_2024_FLOORS):
        swap_picture(s9, FIG_2024_FLOORS)
    add_note(s9, "Comment responses: (1) examples filled in; for the 0.20 floor "
             "the requirements are 5 of 10 or 30 of 100. (2) Yes - raw "
             "selection always carried MIN_TRAIN_FLAGGED = 10. (3) The old "
             "chart was scored on 2023 (its rules trained on 2022+2024). (4) "
             "Replaced with the 2024-scored version: rules re-mined on 2022-23 "
             "(methods/floor_definitions_2024_figure.R). Same story on 2024: "
             "raw floors overpromise; LCB floors deliver at or above their "
             "number (0.30 floor -> 0.336, 0.40 -> 0.475).")

    # C11E (tuning slide): engine-combo chart in; old chart -> appendix
    s12 = find_slide(p, "Tuning: what moved held-out performance")
    swap_picture(s12, "methods/compare_engines_v2/combo_sweeps.png")
    add_note(s12, "Comment 'replace with model performance rpart vs xgboost vs "
             "two + ranger': chart swapped to methods/compare_engines_v2/"
             "combo_sweeps.png (pairs vs singles, xgboost+ranger highlighted). "
             "The one-at-a-time xgboost sweep moved to an appendix slide at "
             "the end.")
    s_app = p.slides.add_slide(s12.slide_layout)
    t_app = copy_shape(s12, s_app, "Tuning: what moved")
    set_para(t_app.text_frame, "Tuning: what moved",
             "Appendix: one-at-a-time xgboost sweep")
    s_app.shapes.add_picture("methods/parameter_tuning_v2/v2_tuning_xgboost.png",
                             Inches(2.8), Inches(0.85), width=Inches(4.4))
    add_note(s_app, "Created for the comment on the tuning slide: the old "
             "figure, preserved in an appendix.")
    for fig, ttl in (
            ("methods/parameter_tuning_v2/v2_tuning_ranger.png",
             "Appendix: one-at-a-time ranger sweep"),
            ("methods/parameter_tuning_v2/v2_lcbz_sweep.png",
             "Appendix: filter stringency sweep"),
            ("methods/parameter_tuning_v2/v2_subsample_fine.png",
             "Appendix: subsampling sweep (claim retired on 2024)"),
            ("methods/compare_engines_v2/engine_sweeps.png",
             "Appendix: single engines vs pairs")):
        s_x = p.slides.add_slide(s12.slide_layout)
        t_x = copy_shape(s12, s_x, "Tuning: what moved")
        set_para(t_x.text_frame, "Tuning: what moved", ttl)
        s_x.shapes.add_picture(fig, Inches(2.8), Inches(0.85), width=Inches(4.4))
        add_note(s_x, "User request: existing sweep figure added to the appendix.")

    # user request: pre-era (RuleFit/rpart) figures on strata + separate
    # elderly/disabled model, placed BEFORE the xgboost+ranger engine charts
    s11 = find_slide(p, "More strata seemed to be a big win")
    s11.shapes.add_picture(
        "compare_models_by_HHsize_vs_pooled/earn_inc_pr_overall.png",
        Inches(0.44), Inches(1.85), width=Inches(4.3))
    s11.shapes.add_picture(
        "compare_models_by_HHsize_vs_pooled/optimal_HH_split_test_separate_ESAP_model_pr_overall_schemes.png",
        Inches(5.15), Inches(1.85), width=Inches(4.3))
    add_note(s11, "User request: the pre-era (RuleFit/rpart) figures. Left: "
             "stratified-by-HH-size vs pooled (the +47% precision era). Right: "
             "separate elderly/disabled (ESAP) model on top of the HH split - "
             "also a win on the rpart stack. Both from "
             "compare_models_by_HHsize_vs_pooled/; the v2-engine slides that "
             "follow show both advantages shrinking once mtry-2 ensembles get "
             "the stratifiers as features.")

    # C121 (strata slide): 2024 re-run + measured attribution
    s14 = find_slide(p, "Split by coarse household-size strata")
    tf14 = find_tf(s14, "Explicit HH size strata")
    set_para(tf14, "Explicit HH size strata",
             "Explicit HH-size strata mattered a lot on the earlier CART/"
             "RuleFit stack (+47% precision). On the v2 ensembles, pooling "
             "matches the split on precision and gets ~90% of its reach.")
    st = strata_numbers()
    if st:
        pooled, coarse = st.get("Pooled (no split)"), st.get("1 / 2-3 / 4+")
        if pooled and coarse:
            add_para_after(
                tf14, "coverage parity",
                "Re-tested on 2024: same picture - pooled vs 3-way precision "
                f"{pooled['mean_precision']} vs {coarse['mean_precision']}, "
                f"reach at the 0.20 floor {pooled['recall_at_020']} vs "
                f"{coarse['recall_at_020']}.")
        fig = os.path.join(STRATA, "strata_sweeps.png")
        if os.path.exists(fig):
            swap_picture(s14, fig)
    add_note(s14, "Comment 'Didn't we redo this for 2024?': we had not (the "
             "year-swap covered engines/stringency/subsample/menu, not "
             "strata). Re-run 2026-07-12: methods/run_strata_yearswap.R -> "
             "methods/compare_hh_strata_v2/yearswap_train2223_test24/. "
             "Numbers added to the slide; figure swapped to the 2024-scored "
             "version. Also corrected the engine attribution: the measured "
             "contrast is pre-era CART/RuleFit vs the v2 ensembles - there is "
             "no xgboost-only strata test.")

    # C122 (state-scale slide): review + mechanism
    s17 = find_slide(p, "At state scale, confidence bounds alone")
    tf17 = find_tf(s17, "Requiring each rule to flag")
    set_para(tf17, "Requiring each rule to flag",
             "- Requiring >= 30 flagged training cases (with a raw 0.30 floor) "
             "changed collapse into gentle deflation: median 0.33 train -> "
             "0.21 held-out; ~1% of rules at zero.")
    add_para_after(tf17, "National scale hides this",
                   "Why: at 5-10 cases the bound only passes near-perfect raw "
                   "precision - exactly where luck concentrates.")
    add_note(s17, "Comment 'explain / review': verified against the Virginia "
             "year-split logs (custom_one_off/virginia/): bound-only arm "
             "median 2024 precision 0.000, 262 of 583 rules caught nothing; "
             "support-floor arm 0.326 -> 0.211 (~1/3 deflation), 7 of 522 at "
             "zero. One correction: the second arm used a RAW 0.30 floor with "
             "n >= 30, not the lower-bound filter. Mechanism line added.")

    # C123 (thin-states slide): superseded by blend
    s18 = find_slide(p, "For thin states, same-era neighbor data")
    add_para_after(find_tf(s18, "identical engines and filters"),
                   "identical engines and filters",
                   "Superseded as a recipe: on 2024, neighbor pools trailed the "
                   "national and blended lists (median 0.25 vs 0.26-0.32). What "
                   "survives: match the era; neighbors beat own-state where "
                   "own-state collapses.")
    add_note(s18, "Comment 'worse than simply blending?': yes as a deployment "
             "recipe - 2024 medians: blend 0.324/0.262 (5%/10%), national "
             "frozen 0.294/0.270, NB neighbor transfer 0.256/0.245 (findings "
             "14, 16). The era-match and neighbor-vs-own lessons survive, so a "
             "one-line supersession was added rather than deleting the slide.")

    # C136 (two-regime slide): concise bullets + alphabetical chart
    s19 = find_slide(p, "Two-regime")
    tb = s19.shapes.add_textbox(Inches(0.44), Inches(0.9), Inches(3.7), Inches(3.4))
    tf19 = tb.text_frame
    tf19.word_wrap = True
    tf19.text = "- Default: the national list, applied in confidence order."
    for line in ("- Fallback: the state's own rules, where they win and the "
                 "state's validation confirms it (triangles).",
                 "- Precursor to the BLEND (end of section), which folds both "
                 "into one ranked list and was best."):
        para = tf19.add_paragraph()
        para.text = line
    for para in tf19.paragraphs:
        for r in para.runs:
            r.font.size = Pt(14)
    swap_picture(s19, os.path.join(DEPLOY, "two_regime_best_budget05.png"))
    add_note(s19, "Comment 'add concise bullets... precursor to the blended "
             "approach': three bullets added; chart swapped for the "
             "alphabetical re-render (same 5% chart).")

    # C126 (donor similarity): era-change measured
    s23 = find_slide(p, "Donor-state similarity")
    add_para_after(find_tf(s23, "Neighbor lists change sharply"),
                   "Neighbor lists change sharply",
                   "Measured: top-5 donor lists keep only ~2 of 5 members "
                   "across eras; a third of states keep <= 1.")
    add_note(s23, "Comment 'Is this change between the two eras true?': yes - "
             "mean top-5 overlap across 49 states is 2.24 (fire), 2.31 (NB), "
             "1.96 (IDF) of 5; 31-33% of states keep at most one neighbor "
             "(similarity_*_2017_2019.csv vs _2022_2024.csv).")

    # C12F (handoff slide): blended numbers + PR chart on its own slide
    s31 = find_slide(p, "What a state is handed")
    set_para(find_tf(s31, "What a state is handed"), "What a state is handed",
             "What a state is handed under the blended approach")
    set_para(find_tf(s31, "One ranked list per state"), "One ranked list per state",
             "- One ranked list per state: state + national rules BLENDED on "
             "the 99% bound, core to the budget + buffer to 3x. 'Typically "
             "run' = rules activated before 2024 capacity filled.")
    old_tbl = next(sh for sh in s31.shapes if sh.has_table)
    L, T = old_tbl.left, old_tbl.top
    old_tbl._element.getparent().remove(old_tbl._element)
    rows = blended_numbers()
    by = {(r["target"], r["budget"]): r for r in rows}
    states = sorted({r["target"] for r in rows})
    half = (len(states) + 1) // 2
    hdr = ["state", "shipped (5%)", "run (5%)", "shipped (10%)", "run (10%)"]

    def mk_table(subset, left):
        gt = s31.shapes.add_table(len(subset) + 1, 5, left, T,
                                  Inches(4.5),
                                  Inches(0.24 * (len(subset) + 1))).table
        for j, h in enumerate(hdr):
            gt.cell(0, j).text = h
        for i, t in enumerate(subset, 1):
            r5, r10 = by[(t, "0.05")], by[(t, "0.1")]
            vals = [t, r5["n_shipped"], r5["n_deployed"],
                    r10["n_shipped"], r10["n_deployed"]]
            for j, v in enumerate(vals):
                gt.cell(i, j).text = str(v)
        for i in range(len(subset) + 1):
            gt.rows[i].height = Inches(0.24)
            for j in range(5):
                c = gt.cell(i, j)
                c.margin_top = c.margin_bottom = Inches(0.005)
                for para in c.text_frame.paragraphs:
                    for r in para.runs:
                        r.font.size = Pt(9)
                        r.font.bold = (i == 0)

    mk_table(states[:half], L)
    mk_table(states[half:], Inches(5.15))
    add_note(s31, "Comment 'Aren't these blended now?': yes - table rebuilt "
             "from the blended lists (blended_frozen_results.csv, 99% bound): "
             "shipped = core+buffer of the blend, typically run = activated at "
             "capacity on 2024. Note: per-state blended list FILES are not yet "
             "exported/committed (frozen_lists/ holds the national-pool and "
             "own-pool versions); say the word and I will export them. The 10% "
             "precision-recall chart got its own slide (next).")

    s_pr = p.slides.add_slide(s31.slide_layout)
    t_pr = copy_shape(s31, s_pr, "What a state is handed")
    set_para(t_pr.text_frame, "What a state is handed",
             "Blend vs national vs own rules, 10% budget")
    s_pr.shapes.add_picture(
        os.path.join(DEPLOY, "blend_vs_state_vs_national_budget10.png"),
        Inches(2.9), Inches(0.75), width=Inches(4.3))
    add_note(s_pr, "Comment 'Add chart with precision recall on a different "
             "slide': the 10% per-state chart (the 5% version sits on the "
             "review-budgets slide).")
    move_slide(p, slide_pos(p, s_pr), slide_pos(p, s31) + 1)

    # alphabetical chart convention on remaining state charts
    s20 = find_slide(p, "Freeze each state's list in advance")
    swap_picture(s20, os.path.join(DEPLOY, "frozen_lists_panels_budget05.png"))
    add_note(s20, "Chart re-rendered with states alphabetical (same 5% chart).")
    s21 = find_slide(p, "Evaluate at review budgets")
    swap_picture(s21, os.path.join(DEPLOY, "blend_vs_state_vs_national_budget05.png"))
    add_note(s21, "Chart re-rendered with states alphabetical (same 5% blend chart).")
    s22 = find_slide(p, "The national rules, deployed a year ahead")
    swap_picture(s22, os.path.join(DEPLOY, "deploy_national_dotplot_budget10.png"))
    add_note(s22, "Chart re-rendered with states alphabetical.")

    # stale-number fixes
    s28 = find_slide(p, "The deployed list is a few dozen rules")
    set_para(find_tf(s28, "26-54 rules at these budgets"),
             "26-54 rules at these budgets",
             "- A rule whose cases are all already flagged adds zero and "
             "always 'fits', so admitted counts run 10-20k. The union is built "
             "by the rules adding new cases: 16-39 at 5% and 34-67 at 10% "
             "across 18 states (medians 27 / 52).")
    add_note(s28, "Stale-number fix: 26-54 was the first three example states; "
             "replaced with the measured 18-state range.")

    s33 = find_slide(p, "live on github")
    set_para(find_tf(s33, "Results for each of the states"),
             "Results for each of the states",
             "- Ranked rule lists for each state discussed (~60-150 rules for "
             "a 5% review budget, ~70-290 for 10%)")
    add_note(s33, "Number check: the per-state lists in frozen_lists/ run "
             "62-148 rules at the 5% sizing and 72-286 at 10% (~80-200 clipped "
             "both ends). These are the national-pool frozen lists; the "
             "blended versions are not yet exported as files.")


# ── variant B extras (short, in the user's voice) ─────────────────────────────
def apply_variant_b(p):
    s2 = find_slide(p, "Approach is to extract rules")
    tf2 = find_tf(s2, "Approach is to extract rules")
    set_para(tf2, "Optimizing",
             "What to optimize: household-size strata; elderly/disabled as a "
             "feature vs a stratum")
    set_para(tf2, "Which are best",
             "Which rules serve a state best: national, similar-state, or own "
             "- and how to combine them")
    set_para(tf2, "Fewer rules with less chance",
             "Fewer rules with fewer false positives, or more rules with "
             "stricter selection? More rules + a strict lower bound won.")
    add_note(s2, "revB: tightened the three framing lines; content unchanged.")

    s6 = find_slide(p, "Limitations of the public QC data")
    add_para_after(find_tf(s6, "bounds what any public-data model"),
                   "bounds what any public-data model",
                   "Range: 43% (New Jersey) to 91% (Georgia); national 71%.")
    add_note(s6, "revB: added the range so the table reads as examples.")

    s9 = find_slide(p, "Better way to select rules")
    set_para(find_tf(s9, "Wilson lower bound"), "Wilson lower bound",
             "- Keep a rule only if the 99% Wilson lower bound of its training "
             "precision clears the floor: delivered precision then meets the "
             "floor, with ~50% more recall at 0.20 than raw selection.")
    add_note(s9, "revB: lead bullet rewritten - it merged the calibration and "
             "recall facts into one hard-to-parse sentence.")

    s11 = find_slide(p, "More strata seemed to be a big win")
    tb = s11.shapes.add_textbox(Inches(0.44), Inches(0.82), Inches(9.0), Inches(0.9))
    tf11 = tb.text_frame
    tf11.word_wrap = True
    tf11.text = ("- On the rpart/RuleFit stack, HH-size strata and a separate "
                 "elderly/disabled model both won clearly (below).")
    para = tf11.add_paragraph()
    para.text = ("- On the v2 ensembles both gaps nearly closed - next slides.")
    for para in tf11.paragraphs:
        for r in para.runs:
            r.font.size = Pt(14)
    add_note(s11, "revB: this slide had a title only; added two short set-up "
             "bullets (delete if you meant it as a bare divider).")

    s16 = find_slide(p, "state-generated and nationally-generated")
    add_para_after(find_tf(s16, "state-generated and nationally-generated"),
                   "state-generated and nationally-generated",
                   "(Answer, end of section: rank every rule by its own 99% "
                   "lower bound - the bound prices the certainty difference.)")
    add_note(s16, "revB: added the destination so the section poses a question "
             "it then answers.")

    s25 = find_slide(p, "adaptation does not beat the national order")
    s_ad = p.slides.add_slide(s25.slide_layout)
    t_ad = copy_shape(s25, s_ad, "adaptation does not beat")
    set_para(t_ad.text_frame, "adaptation does not beat",
             "The adaptation schemes, drawn")
    s_ad.shapes.add_picture(os.path.join(DEPLOY, "adaptation_arms_budget10.png"),
                            Inches(2.9), Inches(0.75), width=Inches(4.3))
    add_note(s_ad, "revB: the adaptation tables had no figure; this is "
             "methods/visualize_state_adaptation_v2.R (filled = national "
             "as-is; open = adapted; NJ/MS/DC = adaptation wins, MA/MI/WA = "
             "losses).")
    move_slide(p, slide_pos(p, s_ad),
               slide_pos(p, find_slide(
                   p, "10% review budget: national as-is vs adapted")) + 1)


def build(variant, out_path):
    shutil.copy(SRC, out_path)
    p = Presentation(out_path)
    apply_variant_a(p)
    if variant == "B":
        apply_variant_b(p)
    p.save(out_path)
    print(f"built {out_path} with {len(p.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build("A", OUT_A)
    build("B", OUT_B)
