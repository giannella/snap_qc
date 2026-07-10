# Builds lessons_getting_more_signal_from_snap_qc_data.pptx from the how_to
# deck (theme + text-box formatting cloned from its appendix slides).
#
# Framing rules (2026-07-08 review): modeling lessons only, each backed by a
# number or chart from our own runs; no claims about real-world phenomena or
# data semantics (the audience knows SNAP data better than we do); no
# over-generalized slogans; deduction/'other' errors are NOT played up (most
# states treat them as low-value).
import copy
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt

SRC = "how_to_use_the_snap_qc_inclusion_rules_scripts.pptx"
OUT = "lessons_getting_more_signal_from_snap_qc_data.pptx"
TEMPLATE_IDX = 16          # "Appendix: lucky rules, measured" - title/body/tag boxes
TAG = "SNAP QC modeling lessons - July 2026"

shutil.copy(SRC, OUT)
p = Presentation(OUT)
tmpl = p.slides[TEMPLATE_IDX]

shapes = [sh for sh in tmpl.shapes if sh.has_text_frame]
tmpl_title = next(sh for sh in shapes if sh.text_frame.text.startswith("Appendix:"))
tmpl_body  = next(sh for sh in shapes if sh.text_frame.text.startswith("- "))
tmpl_tag   = next(sh for sh in shapes if "pipeline evidence" in sh.text_frame.text)


def set_line(tf, text):
    first = tf.paragraphs[0]
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    if first.runs:
        first.runs[0].text = text
        for r in first.runs[1:]:
            r.text = ""
    else:
        first.text = text


def set_body(tf, lines):
    set_line(tf, lines[0])
    first_p = tf.paragraphs[0]._p
    for line in lines[1:]:
        new_p = copy.deepcopy(first_p)
        first_p.addnext(new_p)
        first_p = new_p
    for para, line in zip(tf.paragraphs, lines):
        if para.runs:
            para.runs[0].text = line
            for r in para.runs[1:]:
                r.text = ""


def add_lesson(title, bullets, figure=None, fig_aspect=None, table=None,
               table_font=12, table_row_h=0.32, table_w=6.5):
    s = p.slides.add_slide(tmpl.slide_layout)
    for src_sh in (tmpl_title, tmpl_body, tmpl_tag):
        s._element.spTree.append(copy.deepcopy(src_sh._element))
    boxes = [sh for sh in s.shapes if sh.has_text_frame]
    title_sh = next(sh for sh in boxes if sh.text_frame.text.startswith("Appendix:"))
    body_sh  = next(sh for sh in boxes if sh.text_frame.text.startswith("- "))
    tag_sh   = next(sh for sh in boxes if "pipeline evidence" in sh.text_frame.text)
    set_line(title_sh.text_frame, title)
    set_body(body_sh.text_frame, bullets)
    set_line(tag_sh.text_frame, TAG)
    wrapped = sum(max(1, -(-len(b) // 105)) for b in bullets)
    fig_top = 0.85 + 0.28 * wrapped + 0.2
    if figure:
        aspect = fig_aspect if fig_aspect else 5 / 8
        max_h = 5.05 - fig_top
        h = min(max_h, 3.4)
        w = h / aspect
        if w > 9.2:
            w = 9.2
            h = w * aspect
        s.shapes.add_picture(figure, Inches((10 - w) / 2), Inches(fig_top),
                             width=Inches(w), height=Inches(h))
    if table:
        rows, cols = len(table), len(table[0])
        tw = Inches(table_w)
        th = Inches(table_row_h * rows)
        gt = s.shapes.add_table(rows, cols, Inches((10 - table_w) / 2),
                                Inches(fig_top), tw, th).table
        for i, row in enumerate(table):
            gt.rows[i].height = Inches(table_row_h)
            for j, val in enumerate(row):
                c = gt.cell(i, j)
                c.text = val
                c.margin_top = c.margin_bottom = Inches(0.01)
                for para in c.text_frame.paragraphs:
                    for r in para.runs:
                        r.font.size = Pt(table_font)
                        r.font.bold = (i == 0)
    return s


n_before = len(p.slides._sldIdLst)

# ── 1. title ─────────────────────────────────────────────────────────────────
s = p.slides.add_slide(tmpl.slide_layout)
for src_sh in (tmpl_title, tmpl_body):
    s._element.spTree.append(copy.deepcopy(src_sh._element))
boxes = [sh for sh in s.shapes if sh.has_text_frame]
t_sh = next(sh for sh in boxes if sh.text_frame.text.startswith("Appendix:"))
b_sh = next(sh for sh in boxes if sh.text_frame.text.startswith("- "))
set_line(t_sh.text_frame, "Getting more signal out of SNAP QC data: modeling lessons")
set_body(b_sh.text_frame, [
    "What we learned from building, tuning, and stress-testing an error-mining pipeline on the public QC files.",
    "Every claim below is backed by a measurement from our own runs - configurations tried head-to-head on held-out years.",
    "Code and full evidence: https://github.com/giannella/snap_qc (methods/modeling_findings.md)"])

# ── 2. scope ─────────────────────────────────────────────────────────────────
add_lesson(
    "What we ran",
    ["- Task: mine interpretable rules that flag high-error-risk cases; judge every configuration on a year of data it never saw (train 2022+2024, test 2023).",
     "- Compared head-to-head: tree engines and their pairings, ensemble sizes, learning rates, subsampling, tree depth, split randomness, stratification schemes, selection statistics, filter stringencies, and state-level deployment strategies.",
     "- The lessons below are the settings and practices that moved held-out performance - and the ones that didn't."])

# ── 3-4. data-build effects (measured) ───────────────────────────────────────
add_lesson(
    "Two data-build choices changed results more than any tuning",
    ["- Cases whose review found a second error element did not fit our single-error reconstruction and were being dropped: 31% of all above-threshold error cases. Keeping them (with a flag) raised qualifying rules from 3,741 to 11,018 at the same filter - more error data tightens every statistical bound.",
     "- The fix did not rewrite the rules' content: 93% of the previous rule set survived the rebuild - 63% as the same variables and directions with shifted cutpoints, 29% matched by a new rule flagging a highly overlapping case set (Jaccard >= 0.5 on the test year), 1% exact; only 7% dropped with no counterpart.",
     "- Nor did the restored cases add a new pattern type: brand-new rules' catches were 34% restored-case vs a 32% base - the gain was statistical power, not different signal.",
     "- Rows with blank optional-deduction fields were also being dropped: ~16% of Washington's caseload. Zero-filling with an imputed-flag kept them.",
     "- Both defects were invisible in model metrics and found only by reconciling row and error counts against the raw files, state by state. We now diff mined-rule sets across data builds (exact / shifted / overlapping / dropped / new) to audit any data change."])

add_lesson(
    "Check how much of the error population your data can see",
    ["- We reconciled the public files against the QC technical documentation's exclusion counts, per state and year (2022-24).",
     "- The public files' share of each state's error cases ranges from 43% to 81% - the excluded cases are the ineligible determinations.",
     "- This bounds what any public-data model can deliver per state; we quote it alongside every state result."],
    table=[["state", "visible share of error cases"],
           ["New Jersey", "43%"],
           ["Tennessee", "51%"],
           ["Arkansas / Missouri / Utah", "~53%"],
           ["Washington / Virginia / Louisiana", "78-81%"],
           ["national", "71%"]])

# ── 5-8. selection + tuning lessons ──────────────────────────────────────────
add_lesson(
    "Selecting rules on raw training precision delivered half of it",
    ["- Our 'train precision >= 0.20' shortlist delivered ~0.10 on the held-out year. Before any selection, train precision was nearly unbiased for holdout (median gap -0.003, r = 0.83); rules selected on the HOLDOUT >= 0.20 showed median train precision 0.116.",
     "- So the decay is selection noise (regression to the mean), not overfitting or year drift - the same rules held ~3.9x lift on 2018-19 vs ~3.5x on 2023.",
     "- Any 'generate many candidates, keep the best-measured' step has this problem, whatever the model class."],
    figure="presentation_figures/winners_curse_raw_vs_lcb.png", fig_aspect=1.0)

add_lesson(
    "Filter on a lower confidence bound of training precision",
    ["- Keep a rule only if the one-sided 99% Wilson lower bound of its training precision clears the floor. At matched delivered precision (~0.20), lower-bound selection caught 12.8% of errors vs 8.2% for raw-precision selection.",
     "- Same rule pool, measured three ways: raw floors overpromise even after a lower-bound gate (floor 0.50 delivers 0.34); lower-bound floors deliver at or above their number (floor 0.30 delivers 0.38, floor 0.40 delivers 0.51).",
     "- The lower bound is the only selection statistic we tested whose value survives contact with a new year."],
    figure="presentation_figures/floor_definitions_educational.png", fig_aspect=9.5 / 13)

add_lesson(
    "Bigger ensembles widen the menu, not the frontier",
    ["- 1000-round/1000-tree ensembles traced the SAME precision-recall frontier as small ones - but produced several times more distinct rules clearing any given floor (~2.6-5x the filtered inventory).",
     "- That inventory is what states use: reviewers veto rules they distrust and need substitutes.",
     "- The stringent lower-bound filter is what keeps large pools safe: it removes the extra selection noise that more candidates would otherwise inject."],
    figure="presentation_figures/mine_big_filter_stringently.png", fig_aspect=5 / 8)

add_lesson(
    "Tuning: what moved held-out performance and what didn't",
    ["- Engine PAIRING mattered: xgboost + random forest (ranger) beat either engine alone and beat bagged-CART + ranger, everything else held equal.",
     "- Within engines: mtry = 2 beat mtry = 1 across the frontier; slow learning (eta 0.02, 1000 rounds) beat fast; depth 4 sufficed; subsample barely matters (see the replication on the next slide).",
     "- Most other knobs did not move the frontier - we verified by varying one setting at a time around the final configuration."],
    figure="methods/parameter_tuning_v2/v2_tuning_xgboost.png", fig_aspect=2100 / 2700)

add_lesson(
    "Re-test your selection choices on a year that never judged them",
    ["- Every configuration choice above was judged on the same held-out year (2023), so the selection PROCEDURE itself could be overfit to one year. We re-ran the decisive comparisons with the year roles swapped (train 2022+2023, test 2024), with expected outcomes and failure criteria written down BEFORE the run.",
     "- Three of four claims replicated: the engine pairing still leads recall (margin thinner, 2.1pp vs 3pp+), stringency still buys precision monotonically, big ensembles still widen the menu ~7x at ~0.02 precision cost.",
     "- One claim failed and is retired: 'low subsampling beats high' - on 2024 all nine settings from 0.15 to 0.80 land within 0.005 (right panel). The retraction is what makes the survivors credible: the procedure demonstrably can say no."],
    figure="presentation_figures/replication_selection_studies.png", fig_aspect=4.8 / 12)

# ── 9-10. scoring + stratification ───────────────────────────────────────────
add_lesson(
    "Score rules against all error types, not just the mined type",
    ["- Rules mined for one error type also flag cases carrying other types. Scored only against the mined type, our earned-income rules measured 8% precision on the held-out year; scored against any above-threshold error on the same flags, 19%.",
     "- The gap was similar across frames (~2x). Reporting only the narrow metric understates what reviewers would actually find.",
     "- We now compute both in every run (dashed vs solid below)."],
    figure="inclusion_rules_by_hh_size_v2/earned_income_lcb_sweep.png", fig_aspect=0.5)

add_lesson(
    "Split by coarse household-size strata (1 / 2-3 / 4+)",
    ["- Household-size strata 1 / 2-3 / 4+ beat pooled modeling on recall reach and filtered inventory (~5x); a 5-way split was worse on the same test year.",
     "- Adding an elderly/disabled STRATUM on top bought nothing: an indicator variable achieved the same held-out frontier and the same coverage of those households - we checked coverage parity explicitly.",
     "- We test 'new stratum vs new variable' empirically before adding strata; strata cost data and the variable usually suffices."],
    figure="methods/compare_hh_strata_v2/strata_sweeps.png", fig_aspect=1650 / 2700)

# ── 11-12. small samples + transfer ──────────────────────────────────────────
add_lesson(
    "At state scale, confidence bounds alone were not enough",
    ["- Mining directly on one state's data (~2,000-3,000 cases/year), rules passing the same 99% lower-bound filter with small support had median HOLDOUT precision of zero at support >= 5.",
     "- Requiring each rule to flag >= 30 training cases changed the failure mode from collapse to gentle deflation (~1/3 below training estimates).",
     "- National scale hides this: with 100k+ cases the bound rarely admits tiny-support rules in the first place."])

add_lesson(
    "For thin states, same-era neighbor data beat more own-state data",
    ["- Louisiana's own 2022-24 mining collapsed on holdout. Adding its 2017-19 data made it WORSE (median holdout precision fell to zero).",
     "- Training on five similar states' 2022-24 data - never touching Louisiana - delivered 14% precision at 49% of error dollars on Louisiana 2022-24.",
     "- Both comparisons used identical engines and filters; the only differences were WHERE and WHEN the training data came from."])

add_lesson(
    "Evaluate at review budgets, not just statistical floors",
    ["- Confidence floors produce whatever workload they produce: on the rebuilt data, the 0.20 floor's rule union flags ~half the caseload (which is how 'catch 79% of error dollars' happens). States plan around review capacity, not floors.",
     "- So we also report budget-filled performance: add rules in descending lower-bound order until 5% or 10% of the caseload is flagged. In the deployment test (trained on 2022-23 only, scored on each state's 2024), the national rules deliver median 0.30 precision catching 16% of error dollars at a 5% budget, and 0.27 catching 25% at 10% - against 8-17% base rates.",
     "- The budget lens also corrects floor artifacts: Mississippi's transferred rules stopped firing entirely at a 0.30 floor, yet delivered among the best 5%-budget precision of any state - the rules were fine, the floor wasn't."],
    figure="methods/state_similarity_v2/transfer_benchmark/best_pool_pr_by_budget.png",
    fig_aspect=5.8 / 10)

add_lesson(
    "The national rules, deployed a year ahead, state by state",
    ["- The deployment test: rules mined on all states' 2022-23 data, applied in national-confidence order until 10% of the state's caseload is flagged, scored on the state's 2024 cases - a year no rule ever saw.",
     "- Every state clears its base rate: precision 0.18-0.37 against 8-17% base rates, a 1.5-3.4x lift over random review.",
     "- Including the state's own 2022-23 rows in the mining changes little vs holding the state out entirely (median difference ~0.00; per state -0.08 to +0.06): with a truly unseen test year, the national numbers a state is quoted are already honest."],
    figure="methods/state_similarity_v2/transfer_benchmark_train2223_test24/deploy_national_dotplot_budget10.png",
    fig_aspect=4.4 / 6.6)

add_lesson(
    "Donor-state similarity: tested, and the national pool held up",
    ["- We compute state similarity from the QC microdata itself (which risk patterns fire in a state's caseload; rarely-firing patterns weighted more heavily) and pick each state's 5 nearest donors. Neighbor lists change sharply between 2017-19 and 2022-24, so we compute them per era.",
     "- Scored within the same era, 5-donor pools matched or beat the held-out national pool in 6 of 12 states at a 10% budget (median 0.278 vs 0.245). Re-tested as a deployment (train 2022-23, score 2024), the advantage did not survive: national median 0.27-0.30 vs transfer 0.25-0.26 across budgets.",
     "- Mining a state's own data alone is the high-variance option: the best single results anywhere (Connecticut 0.416, Virginia 0.371 at 10%) but also the worst - Washington's own rules delivered 0.05-0.08 precision on 2024, BELOW its 8.5% base rate. Neighbor pools are the fallback where own-state mining fails, not the default."])

# ── 15-18. adapting national rules to a state (deployment test) ──────────────
add_lesson(
    "Adapting the national rules to a state: three schemes we tested",
    ["- filtered - keep the national rules exactly as they are, but re-qualify each on the state's own 2022-23 data: deploy only rules whose 90% lower-bound precision clears 0.20 with >= 30 flagged cases, applied in state-confidence order.",
     "- tuned - grid-search each rule's thresholds (+/-10-25%) on the state's data; qualify variants at the same bar (90% lower bound >= 0.20, >= 30 flagged); deploy each rule's highest-lower-bound variant - the most statistically defensible version.",
     "- hybrid - the scheme our July grid-search study settled on: qualify variants at 90% lower bound >= 0.20 with >= 5 flagged, then deploy the qualifying variant catching the most error DOLLARS on state training data - a careful gate with an aggressive pick.",
     "- All arms are budget-filled on the state's 2024 cases: rules added in descending confidence order until 5% or 10% of the caseload is flagged. The state's 2024 data is never used to pick rules, only to count flags.",
     "- Note: tuning arms search the top 500 national rules by confidence; the national and filtered arms draw from the full ~45k-rule pool."])

import csv as _csv


def _adapt_rows():
    path = ("methods/state_similarity_v2/transfer_benchmark_train2223_test24/"
            "deployment_state_adaptation.csv")
    with open(path, newline="") as f:
        return list(_csv.DictReader(f))


_ADAPT = _adapt_rows()
_ARMS = ["national_asis", "filtered", "tuned", "hybrid"]
_ARM_LABELS = ["national as-is", "filtered", "tuned", "hybrid"]


def _cell(r):
    if not r or r["precision"] in ("", "NA"):
        return "-"
    return "%.3f @ %.0f%% ($%.0f%%)" % (float(r["precision"]),
                                        100 * float(r["recall"]),
                                        100 * float(r["dollar_recall"]))


def _adapt_table(budget):
    rows = {a: {r["target"]: r for r in _ADAPT
                if r["approach"] == a and float(r["budget"]) == budget}
            for a in _ARMS}
    order = sorted(rows["national_asis"],
                   key=lambda t: -float(rows["national_asis"][t]["precision"]))
    tbl = [["state", "rules qualified (f/t/h)"] + _ARM_LABELS]
    for t in order:
        nq = "/".join(rows[a][t]["n_deployable_rules"]
                      for a in ("filtered", "tuned", "hybrid"))
        tbl.append([t, nq] + [_cell(rows[a].get(t)) for a in _ARMS])
    return tbl


def _med(vals):
    v = sorted(vals)
    n = len(v)
    return (v[n // 2] if n % 2 else (v[n // 2 - 1] + v[n // 2]) / 2)


def _med_cell(arm, budget):
    rs = [r for r in _ADAPT if r["approach"] == arm and float(r["budget"]) == budget]
    return "%.3f / %.0f%%" % (_med([float(r["precision"]) for r in rs]),
                              100 * _med([float(r["dollar_recall"]) for r in rs]))


add_lesson(
    "Median across 18 states: adaptation does not beat the national order",
    ["- No adaptation scheme beats deploying the national list in national-confidence order on the median state; national as-is takes the most per-state precision wins at both budgets (9 of 18 at 5%, 7 of 18 at 10%).",
     "- Adaptation pays where the national list underperforms (New Jersey 0.10 -> 0.27, Mississippi 0.25 -> 0.42 at 5%) and hurts where it is already strong (Michigan 0.40 -> 0.19, Washington 0.30 -> 0.08 at 5%): small-sample qualification re-introduces the selection noise the national bound removed.",
     "- The tuned arm is too conservative to deploy: 1-125 rules survive per state and 11 of its 36 runs cannot fill even half their budget. Cells: median precision / median share of error dollars, on 2024."],
    table=[["scheme", "5% budget", "10% budget"]] +
          [[lab, _med_cell(a, 0.05), _med_cell(a, 0.10)]
           for a, lab in zip(_ARMS, _ARM_LABELS)])

add_lesson(
    "State by state, 5% review budget: national as-is vs adapted",
    ["- precision @ recall (share of error dollars) on the state's 2024 cases, trained on 2022-23 only. 'Rules qualified' counts RULES from the national pool clearing the state's bar (rules overlap heavily, so counts can exceed the state's caseload; a high base rate makes the 0.20 bar easy to clear). Tuning arms search the top 500 national rules by confidence; the national and filtered arms draw from the full ~45k-rule pool."],
    table=_adapt_table(0.05), table_font=9, table_row_h=0.23, table_w=9.2)

add_lesson(
    "State by state, 10% review budget: national as-is vs adapted",
    ["- precision @ recall (share of error dollars) on the state's 2024 cases, trained on 2022-23 only. 'Rules qualified' counts RULES from the national pool clearing the state's bar (rules overlap heavily, so counts can exceed the state's caseload; a high base rate makes the 0.20 bar easy to clear). Tuning arms search the top 500 national rules by confidence; the national and filtered arms draw from the full ~45k-rule pool."],
    table=_adapt_table(0.10), table_font=9, table_row_h=0.23, table_w=9.2)

def _contrib_table():
    path = ("methods/state_similarity_v2/transfer_benchmark_train2223_test24/"
            "contributing_rules_summary.csv")
    with open(path, newline="") as f:
        rows = list(_csv.DictReader(f))
    tbl = [["state", "budget", "review cap (cases)", "rules admitted by the scan",
            "rules that build the union"]]
    for r in rows:
        tbl.append([r["target"], "%.0f%%" % (100 * float(r["budget"])),
                    r["cap_cases"], r["rules_counted_used"], r["rules_contributing"]])
    return tbl


add_lesson(
    "The deployed list is a few dozen rules, not thousands",
    ["- The budget fill scans the whole ~45k national pool in confidence order and admits any rule whose newly flagged cases still fit under the review cap - it does not stop at a rule count.",
     "- A rule whose cases are all already flagged adds zero and always 'fits', so admitted-rule counts run 10-20k. The union is actually built by the rules adding at least one new case: 26-54 rules at these budgets. Deploying exactly those reproduces the identical flagged set.",
     "- Which rules those are depends only on the ranked list (trained on 2022-23) and the state's incoming caseload - no error outcomes - so a state reproduces the short list the moment it applies the ranking to its own cases."],
    table=_contrib_table(), table_font=11, table_row_h=0.28, table_w=8.5)

add_lesson(
    "The deliverable: one ranked rule list per state, walked to capacity",
    ["- Built in advance from public data alone: fill the national pool (mined on all states' 2022-23) against the state's own 2022-23 caseload, in confidence order, until the review budget is reached - then keep filling to 3x that depth. The extra rules are ranked BUFFER rules, shipped in the same list.",
     "- To use it, the state activates rules in order, checking only its own flag counts: keep adding while reviews fit capacity. No error outcomes are needed at any step.",
     "- The buffer is what makes the workload land on budget as patterns change: a state never stops reviewing because a list ran dry, and an over-firing list trims the same way. Raw core lists drifted to 2.3-12% of caseload on 2024; the walked lists land at 4.9-5.0% and 9.3-10.0% in all 18 states.",
     "- Median rules actually activated: 23 at a 5% budget, 42 at 10%. Validation recipe: we freeze on public data through the latest year; the state validates on its own newer internal data before relying on it."])

add_lesson(
    "Freeze each state's list in advance: it costs almost nothing",
    ["- The frozen list (built on 2022-23, walked to capacity on 2024) against the full national pool applied at once to the realized 2024 caseload - both at identical review volume.",
     "- Median precision 0.294 vs 0.301 (5% budget) and 0.270 vs 0.275 (10%); error dollars 12% vs 16% and 25% vs 25%. All 18 states clear their base rate."],
    figure="methods/state_similarity_v2/transfer_benchmark_train2223_test24/frozen_lists_panels_budget10.png",
    fig_aspect=5.8 / 6.6)

# ── 21. takeaways ────────────────────────────────────────────────────────────
add_lesson(
    "Takeaways",
    ["- Reconcile your build against the raw files before tuning anything: our two data-build fixes moved results more than any hyperparameter.",
     "- Select on a lower confidence bound, judge on a held-out year, and add hard support floors when samples are small.",
     "- Score against all error types; stratify only where a variable can't do the job; prefer same-era similar data over more mixed data.",
     "- Numbers, charts, and code: github.com/giannella/snap_qc"])

# ── drop the original how_to slides ──────────────────────────────────────────
sldIdLst = p.slides._sldIdLst
for sld in list(sldIdLst)[:n_before]:
    sldIdLst.remove(sld)

p.save(OUT)
print("built", OUT, "with", len(p.slides._sldIdLst), "slides")
