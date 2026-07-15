# RevC: the senior-statistician pass over revB (2026-07-13). His reading:
# the deck's spine is SELECTION INFERENCE - name it, organize around it,
# and say plainly how sure we are of each number. Changes (additions only;
# every change explained in a [RevC - statistician] speaker note):
#   1. "The evaluation ladder" slide closing Part 2 - every number in the
#      talk placed on a rung; each rung harder to fool (new diagram).
#   2. "State rankings are noise..." slide after the deployment dotplot -
#      the CIs overlap, so rank order is not the claim; clearing the base
#      rate everywhere is.
#   3. "What we still don't know" slide before the FY25/26 validation
#      recipe - candid limitations, each paired with what would answer it.
#   4. Sharpened context in notes: why z = 99% (selection multiplicity),
#      why budgets (floors do not control workload), guarantee reading of
#      the LCB ("at least this").
#
#   python methods/build_revC_statistician.py
import copy
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.text.text import _Paragraph

SRC = "lessons_getting_more_signal_from_snap_qc_data_revB.pptx"
OUT = "lessons_getting_more_signal_from_snap_qc_data_revC.pptx"
DEPLOY = "methods/state_similarity_v2/transfer_benchmark_train2223_test24"


def find_slide(p, substr):
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return s
    raise KeyError(substr)


def add_note(s, text):
    tf = s.notes_slide.notes_text_frame
    ex = tf.text
    tf.text = (ex + "\n\n" if ex.strip() else "") + \
        "[RevC - statistician] " + text


def slide_index(p, s):
    for i, x in enumerate(p.slides):
        if x == s:
            return i
    raise KeyError


def move_before(p, s, target):
    lst = p.slides._sldIdLst
    el = list(lst)[slide_index(p, s)]
    lst.remove(el)
    lst.insert(slide_index(p, target), el)


def content_slide(p, donor_anchor, title, bullets, figure=None,
                  fig_box=None, body_box=None):
    """Clone the donor's title/body/tag shapes onto a new slide."""
    donor = find_slide(p, donor_anchor)
    s = p.slides.add_slide(donor.slide_layout)
    tsh = bsh = None
    for sh in donor.shapes:
        if not sh.has_text_frame:
            continue
        el = copy.deepcopy(sh._element)
        s.shapes._spTree.append(el)
    # identify cloned title (contains donor_anchor) and the longest body
    for sh in s.shapes:
        if sh.has_text_frame and donor_anchor in sh.text_frame.text:
            tsh = sh
    cands = [sh for sh in s.shapes
             if sh.has_text_frame and sh is not tsh and
             len(sh.text_frame.text) > 60]
    if cands:
        bsh = max(cands, key=lambda sh: len(sh.text_frame.text))
    tp = tsh.text_frame.paragraphs[0]
    if tp.runs:
        tp.runs[0].text = title
        for r in tp.runs[1:]:
            r.text = ""
    if bullets and bsh is not None:
        tf = bsh.text_frame
        first = tf.paragraphs[0]
        for para in list(tf.paragraphs[1:]):
            para._p.getparent().remove(para._p)
        anchor = first._p
        for _ in bullets[1:]:
            np_ = copy.deepcopy(first._p)
            anchor.addnext(np_)
            anchor = np_
        for para, line in zip(tf.paragraphs, bullets):
            pp = _Paragraph(para._p, para._parent)
            if pp.runs:
                pp.runs[0].text = line
                for r in pp.runs[1:]:
                    r.text = ""
            else:
                pp.text = line
            for r in pp.runs:
                r.font.size = Pt(13)
            pp.line_spacing = 1.0
            pp.space_after = Pt(5)
        tf.word_wrap = True
        if body_box:
            bsh.left, bsh.top, bsh.width, bsh.height = \
                [Inches(v) for v in body_box]
    elif bsh is not None:
        # no bullets wanted: blank the body
        tf = bsh.text_frame
        for para in list(tf.paragraphs[1:]):
            para._p.getparent().remove(para._p)
        if tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = ""
            for r in tf.paragraphs[0].runs[1:]:
                r.text = ""
    if figure:
        L, T, W = fig_box
        s.shapes.add_picture(figure, Inches(L), Inches(T), width=Inches(W))
    return s


shutil.copy(SRC, OUT)
p = Presentation(OUT)

# 1. the evaluation ladder, closing Part 2
s_lad = content_slide(
    p, "At state scale, confidence bounds alone",
    "Every number here sits on an evaluation ladder",
    ["- Each rung is harder to fool than the one below; the talk climbs the "
     "ladder and reports where each claim stands."],
    figure="presentation_figures/evaluation_ladder.png",
    fig_box=(0.35, 1.55, 9.3))
add_note(s_lad, "New slide. The deck's recurring move - distrust the number, "
         "then re-earn it on a harder test - was implicit; this names it "
         "once so every later slide can locate itself. The ladder also "
         "explains why the FY25/26 ask (final slides) is not a formality: "
         "it is the top rung, and the only one that sees the full error "
         "population.")
move_before(p, s_lad, find_slide(p, "Re-test your selection choices"))

# 2. rankings-are-noise slide, after the deployment dotplot
s_rank = content_slide(
    p, "At state scale, confidence bounds alone",
    "Read the intervals, not the rankings",
    ["- The 95% intervals overlap for most pairs of states: the ORDERING "
     "of states is mostly noise at these sample sizes.",
     "- The durable claim is different: every state's interval clears its "
     "base rate - the lift survives uncertainty, the ranking does not.",
     "- Same reason we do not re-rank rules on small state samples: "
     "selecting on a noisy ranking is the winner's curse again."],
    figure=DEPLOY + "/workshop_national_dotplot_budget10.png",
    fig_box=(4.85, 1.15, 4.9),
    body_box=(0.44, 1.15, 4.2, 3.6))
add_note(s_rank, "New slide, no new analysis - the same 18-state chart, read "
         "the way a statistician would: interval overlap means rank order "
         "is weak evidence, base-rate clearance is strong evidence. This "
         "also inoculates the audience against over-reading their own "
         "state's position.")
move_before(p, s_rank, find_slide(p, "Donor-state similarity"))

# 3. what we still don't know, before the validation recipe
s_unk = content_slide(
    p, "At state scale, confidence bounds alone",
    "What we still don't know",
    ["- Drift beyond one year: our longest honest test is one year ahead; "
     "FY25/26 performance is a forecast, not a measurement.",
     "- The invisible errors: 19-57% of each state's error cases "
     "(ineligible determinations) never enter the public files - the rules "
     "have never seen them.",
     "- The blend's blind spot: a national rule's bound says nothing about "
     "transfer to your state (New Jersey's own rules never enter).",
     "- Engine frontier: constrained forests were the best PERIPHERY miner "
     "we tested, not the best possible one (uniform forests untested).",
     "- Each of these is answerable - mostly by the validation on the next "
     "slide."],
    body_box=(0.44, 1.15, 9.0, 3.8))
add_note(s_unk, "New slide. A limitations slide placed directly before the "
         "validation recipe turns candor into a call to action: three of "
         "the four unknowns are exactly what a state's internal FY25/26 "
         "check answers. Numbers: visibility from findings section 10; NJ "
         "blend case from section 16.")
move_before(p, s_unk, find_slide(p, "Validating on FY25/26"))

# 4. sharpened context in speaker notes on existing slides
add_note(find_slide(p, "Better way to select rules"),
         "Context worth saying aloud: why 99% and not 95%? With 40k+ "
         "candidate rules, the filter faces massive selection multiplicity "
         "- a stringent per-rule bound is the price of mining big. And the "
         "practical reading of the bound: every shortlist number means 'at "
         "least this' - it is the only number in the pipeline that can be "
         "quoted to a state without a discount.")
add_note(find_slide(p, "Evaluate at review budgets"),
         "Context worth saying aloud: floors are statistical statements, "
         "budgets are operational ones. A floor left the workload to "
         "chance (half the caseload); the budget lens fixes the operational "
         "variable and lets precision and dollars be compared apples to "
         "apples across every option in the talk.")
add_note(find_slide(p, "adaptation does not beat the national order"),
         "Sharpen when presenting: this is the deck's second winner's-curse "
         "moment - re-qualifying rules on 2-6k state cases re-introduces "
         "exactly the selection noise the national bound removed. Same "
         "disease, new host; the cure (bigger samples + stringent bounds) "
         "is the same.")

p.save(OUT)
print(f"built {OUT} with {len(p.slides._sldIdLst)} slides")
