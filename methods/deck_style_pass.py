# Deck style pass (2026-07): remove dash-as-punctuation from slide text.
#
# The decks are binary .pptx, so this script IS the reviewable record of the text
# changes. Two kinds of edits, both punctuation/spareness only (no numbers touched):
#   1. em-dashes ("—")  -> colon (almost all are "term - definition").
#   2. " - " hyphen-dash connectors -> period / comma / colon / parentheses per
#      context, since a dash that joins two independent clauses compresses disjoint
#      thoughts. Applied as explicit unique-substring swaps so each is deliberate.
# Plus the two repeated footers, standardized to the title slide's middot style.
#
# Idempotent-ish: re-running simply finds nothing to match. Run-level edits preserve
# bold/italic formatting. Every applied swap is logged; a swap that matches nothing
# is reported so a typo in the map is caught.
#
#   python methods/deck_style_pass.py
import glob
from pptx import Presentation

DASH = "—"

# (old_substring, new_substring). Order matters only if one is a prefix of another.
SWAPS = [
    # --- repeated footers (all decks) ---
    ("SNAP QC modeling lessons - July 2026", "SNAP QC modeling lessons · July 2026"),
    ("Appendix - v2 pipeline evidence", "Appendix: v2 pipeline evidence"),

    # --- lessons deck: content connectors ---
    ("model metrics - found only by reconciling",
     "model metrics. We found them only by reconciling"),
    ("on internal data - it runs in a few min on a 16 GB laptop; a full mining run takes ~45 minutes",
     "on internal data; a full run takes about 45 minutes on a 16 GB laptop"),
    ("On v2 (RF + xbg)", "On v2 (RF + xgboost)"),
    ("flip by year - but the coarse split never loses",
     "flip by year, but the coarse split never loses"),
    (" - on 2024 all nine settings", ". On 2024 all nine settings"),
    ("near-perfect raw precision - exactly where luck concentrates",
     "near-perfect raw precision, exactly where luck concentrates"),
    ("rule - state or national - by its own", "rule (state or national) by its own"),
    ("about transfer - NJ", "about transfer, so NJ"),
    ("at 10% - vs 8-17% base rates", "at 10%, vs 8-17% base rates"),
    ("against ALL error types - a hit is any payment error",
     "against ALL error types. A hit is any payment error"),
    ("public files - the rules have never seen them",
     "public files, so the rules have never seen them"),
    ("errors found, error dollars - at your preferred budget",
     "errors found, error dollars, at your preferred budget"),
    ("by experts - know this", "by experts. Know this"),
    ("realized 2024 caseload - both at identical review volume",
     "realized 2024 caseload, both at identical review volume"),
    ("move the frontier - we verified by varying",
     "move the frontier; we verified this by varying"),
    ("engine tuning - what moved held-out", "engine tuning, what moved held-out"),

    # --- inclusion how-to deck: content connectors ---
    # slide-2 "Key ideas" bold-label paragraphs split the dash across runs, so these
    # target the exact run pieces (bold term run + definition run).
    ("Finding rules - ", "Finding rules: "),
    ("Keeping rules ", "Keeping rules: "),
    ("- a rule tested on only a handful", "a rule tested on only a handful"),
    ("By household size ", "By household size: "),
    ("- cases are split into sizes", "cases are split into sizes"),
    (" - when a flagged case turns out", ": when a flagged case turns out"),
    ("beats either engine alone - ", "beats either engine alone; "),
    ("achieve together - how many cases they flag", "achieve together: how many cases they flag"),
    ("for tuning - under both its adjusted cutoffs and the original national ones - and the whole",
     "for tuning (under both its adjusted cutoffs and the original national ones), and the whole"),
    ("hurt below that - with little data", "hurt below that; with little data"),
    ("both measure 20% - but only the second", "both measure 20%, but only the second"),
    ("delivered accuracy - about 1 in 5 flagged cases having an error - we catch",
     "delivered accuracy (about 1 in 5 flagged cases having an error), we catch"),
    ("the adjustment alone is not enough - also require",
     "the adjustment alone is not enough; also require"),
    ("more accurate - the two curves below", "more accurate; the two curves below"),
    ("tuning detail - the boosted-tree method", "tuning detail for the boosted-tree method"),
    ("tuning detail - the random-forest method", "tuning detail for the random-forest method"),
    ("a single random variable - a little guidance", "a single random variable; a little guidance"),
    ("tuning detail - how much data each tree sees", "tuning detail: how much data each tree sees"),
    ("most of the data (60-80%) - small slices make", "most of the data (60-80%); small slices make"),
    ("find different things - combining both catches", "find different things; combining both catches"),
    ("with other kinds of errors - and those reviews succeed",
     "with other kinds of errors, and those reviews succeed"),
    ("Report the second number - it is what reviewers", "Report the second number: it is what reviewers"),
    ("are the LARGEST error category - 2,007", "are the LARGEST error category: 2,007"),
    ("elderly/disabled households - one model works", "elderly/disabled households: one model works"),

    # --- workshop deck: content connectors (leading "- " bullets left alone) ---
    ("size_v2.R - build a rule list", "size_v2.R: build a rule list"),
    ("gridsearch_v2.R - adjust the rules", "gridsearch_v2.R: adjust the rules"),
    ("snap_qc - everything today", "snap_qc: everything today"),
    ("NATIONAL public QC data - far more", "NATIONAL public QC data: far more"),
    ("(Apache 2.0) - the code", "(Apache 2.0): the code"),
    ("the tuning never saw - not the training data", "the tuning never saw, not the training data"),
    ("Part 1 - what the finder", "Part 1: what the finder"),
    ("Part 2 - tuning rules to one state", "Part 2: tuning rules to one state"),
    ("Part 3 - the fresh results", "Part 3: the fresh results"),
    ("works the same way - swap the variable list", "works the same way; swap the variable list"),
    ("No special hardware - a normal laptop", "No special hardware: a normal laptop"),
    ("HOW MANY rules to deploy - pick the precision floor", "HOW MANY rules to deploy: pick the precision floor"),
    ("columns in your data - this is the main thing", "columns in your data: this is the main thing"),
    ("defaults from our testing - you can run it unchanged", "defaults from our testing, so you can run it unchanged"),
    ("household size - a full national run takes", "household size; a full national run takes"),
    ("_rules_all.csv - every surviving rule", "_rules_all.csv: every surviving rule"),
    ("_rules_highprecision.csv - the shortlist", "_rules_highprecision.csv: the shortlist"),
    ("_lcb_sweep.csv + .png - the menu", "_lcb_sweep.csv + .png: the menu"),
    ("all_frames.csv - shortlists from all error types", "all_frames.csv: shortlists from all error types"),
    ("by confident precision - what do the top rules", "by confident precision: what do the top rules"),
    ("at the flagged cases - count how many were real errors", "at the flagged cases; count how many were real errors"),
    ("MENU, not a mandate - your program experts", "MENU, not a mandate; your program experts"),
    ("never breaks the others - each rule stands on its own", "never breaks the others; each rule stands on its own"),
    ("number in practice - it was set with statistical caution", "number in practice; it was set with statistical caution"),
    ("identically on new data - but the approach we use", "identically on new data, but the approach we use"),
    ("on the test year - which the tuning never touched - under BOTH",
     "on the test year (which the tuning never touched) under BOTH"),
    ("national rules unchanged - the fallback every state has", "national rules unchanged: the fallback every state has"),
    ("name a state - we", "name a state, and we"),
    ("If few rules qualified - that IS the answer", "If few rules qualified, that IS the answer"),
    ("keep the unused rules - they", "keep the unused rules; they"),
    ("reduced settings) - let", "reduced settings). Let"),
    ("couple of hours - same pipeline", "couple of hours; same pipeline"),
    ("numbers or yes/no - recode categories first", "numbers or yes/no; recode categories first"),
    ("public QC files - wage records", "public QC files: wage records"),
    ("evidence-adjusted score - the scripts do this by default", "evidence-adjusted score; the scripts do this by default"),
    ("report both numbers - quote the any-error one", "report both numbers; quote the any-error one"),
    ("snap_qc - scripts, data pipeline", "snap_qc: scripts, data pipeline"),
    ("please reach out - the pipeline improves", "please reach out; the pipeline improves"),
]

def fix_emdash(t):
    t = t.replace(" " + DASH + " ", ": ").replace(DASH, ": ")
    while "  " in t:
        t = t.replace("  ", " ")
    return t

counts = {old: 0 for old, _ in SWAPS}
em = 0
for deck in sorted(glob.glob("slides/*.pptx")):
    p = Presentation(deck)
    touched = False
    for slide in p.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if DASH in run.text:
                        run.text = fix_emdash(run.text); em += 1; touched = True
                    for old, new in SWAPS:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            counts[old] += 1; touched = True
    if touched:
        p.save(deck)

print(f"em-dash runs fixed: {em}")
for old, _ in SWAPS:
    flag = "" if counts[old] else "   <-- MATCHED NOTHING"
    print(f"  {counts[old]:>2}x  {old[:60]!r}{flag}")
