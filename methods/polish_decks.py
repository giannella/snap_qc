# Final polish pass (2026-07-13): professional layout on both rev decks.
#   - traditional title slide (centered title / subtitle / author / date)
#   - styled black-and-white section breaks (revB gets the full 4-part
#     structure; revA keeps its order and gets its existing divider styled)
#   - one NEW optional slide both decks: "Validating on FY25/26" - the title
#     slide promises rules states can test on FY25/26; no slide said how
#   - overflow scan: estimates text height per shape, shrinks fonts that
#     clearly overflow, reports everything it touched
#   - uniform footer tag, sparing bold (section names on the agenda)
# Additions happen BEFORE deletions (part-name collision rule).
#
#   python methods/polish_decks.py
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

REV_A = "lessons_getting_more_signal_from_snap_qc_data_revA.pptx"
REV_B = "lessons_getting_more_signal_from_snap_qc_data_revB.pptx"


def find_slide(p, substr, required=True):
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return s
    if required:
        raise KeyError(substr)
    return None


def find_tf(s, substr):
    for sh in s.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            return sh
    raise KeyError(substr)


def add_note(s, text):
    tf = s.notes_slide.notes_text_frame
    ex = tf.text
    tf.text = (ex + "\n\n" if ex.strip() else "") + "[2026-07-13 polish] " + text


def slide_index(p, s):
    for i, x in enumerate(p.slides):
        if x == s:
            return i
    raise KeyError


def move_before(p, s, target):
    lst = p.slides._sldIdLst
    el = list(lst)[slide_index(p, s)]
    lst.remove(el)
    lst.insert(slide_index(p, target), el)


def delete_slide(p, s):
    lst = p.slides._sldIdLst
    lst.remove(list(lst)[slide_index(p, s)])


def style_runs(tf, size=None, bold=None, italic=None, align=None):
    for para in tf.paragraphs:
        if align is not None:
            para.alignment = align
        for r in para.runs:
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if italic is not None:
                r.font.italic = italic


# ── title slide ───────────────────────────────────────────────────────────────
def polish_title(p, W, H):
    s = p.slides[0]
    title_sh = find_tf(s, "Getting more signal")
    author_sh = find_tf(s, "Eric Giannella")
    tf = title_sh.text_frame
    # single title line + subtitle line
    from pptx.text.text import _Paragraph
    paras = list(tf.paragraphs)
    if paras[0].runs:
        paras[0].runs[0].text = "Getting more signal out of SNAP QC data"
        for r in paras[0].runs[1:]:
            r.text = ""
    subtitle = ("modeling lessons  ·  open-source code  ·  "
                "rules each state can test on FY25/26")
    if len(paras) > 1:
        if paras[1].runs:
            paras[1].runs[0].text = subtitle
            for r in paras[1].runs[1:]:
                r.text = ""
        else:
            paras[1].text = subtitle
        for extra in paras[2:]:
            extra._p.getparent().remove(extra._p)
    title_sh.left, title_sh.top = Inches(0.7), Emu(int(H * 0.30))
    title_sh.width, title_sh.height = Emu(W - Inches(1.4)), Inches(1.7)
    tf.word_wrap = True
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(34)
        r.font.bold = True
    if len(tf.paragraphs) > 1:
        for r in tf.paragraphs[1].runs:
            r.font.size = Pt(16)
            r.font.bold = False
            r.font.italic = False
        tf.paragraphs[1].space_before = Pt(14)
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    author_sh.left, author_sh.top = Inches(0.7), Emu(int(H * 0.74))
    author_sh.width, author_sh.height = Emu(W - Inches(1.4)), Inches(0.8)
    atf = author_sh.text_frame
    if atf.paragraphs[0].runs:
        atf.paragraphs[0].runs[0].text = \
            "Eric Giannella · Georgetown Better Government Lab · July 2026"
        for r in atf.paragraphs[0].runs[1:]:
            r.text = ""
    style_runs(atf, size=14, italic=True, align=PP_ALIGN.CENTER)
    add_note(s, "Title slide restyled: centered traditional layout.")


# ── section break styling ─────────────────────────────────────────────────────
def make_section_break(p, layout_donor, number, title, extra=None):
    s = p.slides.add_slide(layout_donor.slide_layout)
    W, H = p.slide_width, p.slide_height
    box = s.shapes.add_textbox(Inches(0.8), Emu(int(H * 0.34)),
                               Emu(W - Inches(1.6)), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = f"PART {number}"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(13)
        r.font.bold = False
        from pptx.dml.color import RGBColor
        r.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
    para = tf.add_paragraph()
    para.text = title
    para.alignment = PP_ALIGN.CENTER
    para.space_before = Pt(8)
    for r in para.runs:
        r.font.size = Pt(28)
        r.font.bold = True
    if extra:
        para2 = tf.add_paragraph()
        para2.text = extra
        para2.alignment = PP_ALIGN.CENTER
        para2.space_before = Pt(10)
        for r in para2.runs:
            r.font.size = Pt(14)
            r.font.italic = True
    # thin rule above the text
    from pptx.enum.shapes import MSO_SHAPE
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(W * 0.35)),
                            Emu(int(H * 0.30)), Emu(int(W * 0.30)), Pt(1.4))
    ln.fill.solid()
    from pptx.dml.color import RGBColor
    ln.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x40)
    ln.line.fill.background()
    return s


# ── overflow scan ─────────────────────────────────────────────────────────────
def overflow_pass(p, W, H, deck):
    fixed = []
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            tf = sh.text_frame
            # estimate rendered height
            est = 0.0
            for para in tf.paragraphs:
                size = 14.0
                for r in para.runs:
                    if r.font.size is not None:
                        size = r.font.size.pt
                        break
                width_in = max(sh.width / 914400.0, 0.5)
                chars_per_line = max(int(width_in * 150.0 / size * 1.28), 8)
                nlines = max(1, -(-len(para.text) // chars_per_line))
                est += nlines * size * 1.22 / 72.0
            top_in = sh.top / 914400.0
            H_in = H / 914400.0
            if top_in + est > H_in - 0.12 and est > sh.height / 914400.0:
                # shrink fonts one notch and widen box to slide margins
                for para in tf.paragraphs:
                    for r in para.runs:
                        if r.font.size is not None and r.font.size.pt > 11:
                            r.font.size = Pt(max(11, r.font.size.pt - 2))
                tf.word_wrap = True
                fixed.append((deck, i, tf.text[:40]))
            if sh.left + sh.width > W - Inches(0.15):
                sh.width = Emu(W - Inches(0.15) - sh.left)
                tf.word_wrap = True
                fixed.append((deck, i, "width>" + tf.text[:30]))
    return fixed


# ── validation slide (the missing critical piece) ────────────────────────────
def add_validation_slide(p, donor_anchor):
    if find_slide(p, "Validating on FY25/26", required=False) is not None:
        return
    donor = find_slide(p, donor_anchor)
    s = p.slides.add_slide(donor.slide_layout)
    for sh in donor.shapes:
        if sh.has_text_frame:
            el = copy.deepcopy(sh._element)
            s.shapes._spTree.append(el)
    title = find_tf(s, donor_anchor)
    tp = title.text_frame.paragraphs[0]
    if tp.runs:
        tp.runs[0].text = "Validating on FY25/26: the recipe for a state"
        for r in tp.runs[1:]:
            r.text = ""
    body = None
    cands = [sh for sh in s.shapes
             if sh.has_text_frame and len(sh.text_frame.text) > 60 and
             "Validating" not in sh.text_frame.text]
    if cands:
        body = max(cands, key=lambda sh: len(sh.text_frame.text))
    lines = [
        "- The list is frozen on public data through FY24 - before any validation data exists.",
        "- Apply it to internal FY25/26 cases: workload should land near the sized budget.",
        "- The bar: precision comfortably above the state's base error rate.",
        "- If the blend underperforms, the own-pool list is the pre-agreed fallback.",
        "- Public files show only 43-81% of a state's error cases - this internal check is the honest judge.",
    ]
    from pptx.text.text import _Paragraph
    tfb = body.text_frame
    first = tfb.paragraphs[0]
    for para in list(tfb.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    anchor = first._p
    for _ in lines[1:]:
        np_ = copy.deepcopy(first._p)
        anchor.addnext(np_)
        anchor = np_
    for para, line in zip(tfb.paragraphs, lines):
        pp = _Paragraph(para._p, para._parent)
        if pp.runs:
            pp.runs[0].text = line
            for r in pp.runs[1:]:
                r.text = ""
        else:
            pp.text = line
    add_note(s, "NEW optional slide: the title slide promises rules states can "
             "test on FY25/26, but no slide said how. Delete if unwanted.")
    return s


def polish(path, aggressive):
    p = Presentation(path)
    W, H = p.slide_width, p.slide_height
    deck = "B" if aggressive else "A"

    polish_title(p, W, H)

    # additions first (part-name rule) -------------------------------------
    s_val = add_validation_slide(p, "At state scale, confidence bounds alone")
    new_breaks = []
    if aggressive:
        donor = p.slides[2]
        b1 = make_section_break(p, donor, 1, "What the data can and cannot show")
        b2 = make_section_break(p, donor, 2, "Selecting rules that hold up")
        b4 = make_section_break(p, donor, 4,
                                "The deliverable: one blended list per state")
        new_breaks = [b1, b2, b4]

    # find the divider pair; restyle the question one as the Part-3 break
    q_div = find_slide(p, "state-generated and nationally-generated")
    tsh = find_tf(q_div, "Optimizing rule sets for states")
    tp = tsh.text_frame.paragraphs[0]
    if tp.runs:
        tp.runs[0].text = ("PART 3   ·   Optimizing rule sets for states"
                           if aggressive else "Optimizing rule sets for states")
        for r in tp.runs[1:]:
            r.text = ""
    for r in tp.runs:
        r.font.size = Pt(26)
        r.font.bold = True

    # position moves --------------------------------------------------------
    if aggressive:
        move_before(p, new_breaks[0], find_slide(p, "Two data-build choices"))
        move_before(p, new_breaks[1],
                    find_slide(p, "Selecting rules on raw training precision"))
        move_before(p, new_breaks[2],
                    find_slide(p, "Blending state and national rules"))
    if s_val is not None:
        move_before(p, s_val, find_slide(p, "live on github"))

    # deletions last ---------------------------------------------------------
    bare = None
    for sl in p.slides:
        joined = " ".join(sh.text_frame.text for sh in sl.shapes
                          if sh.has_text_frame)
        if "Optimizing rule sets for states" in joined and \
                "state-generated" not in joined:
            bare = sl
            break
    if bare is not None:
        delete_slide(p, bare)

    fixes = overflow_pass(p, W, H, deck)
    p.save(path)
    print(f"polished {path}: {len(p.slides._sldIdLst)} slides; "
          f"overflow fixes: {len(fixes)}")
    for f in fixes:
        print("   ", f)


if __name__ == "__main__":
    polish(REV_A, aggressive=False)
    polish(REV_B, aggressive=True)
