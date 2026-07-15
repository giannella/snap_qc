# Round-3 in-place edits from the fresh revB comments (2026-07-12):
#   129/12A  drop the two per-state adaptation table slides - BOTH decks
#   12B      few-dozen-rules slide: clarify what it is; move after the blend
#            introduction - both decks
#   130      blend intro slide to the top of the blend section - both decks
#   13A      blend PR 10% = second slide of the blend section; add a 5%
#            version to the appendix - both decks
#   135(new) teaching diagram: how the constrained random forest complements
#            xgboost - new slide after "Intuition", both decks
#
#   python methods/apply_deck_revisions_round3.py
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt

REV_A = "lessons_getting_more_signal_from_snap_qc_data_revA.pptx"
REV_B = "lessons_getting_more_signal_from_snap_qc_data_revB.pptx"
DEPLOY = "methods/state_similarity_v2/transfer_benchmark_train2223_test24"
DIAGRAM = "presentation_figures/engines_concept_diagram.png"


def find_slide(p, substr, required=True):
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return s
    if required:
        raise KeyError(f"no slide containing {substr!r}")
    return None


def find_tf(s, substr):
    for sh in s.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            return sh.text_frame
    raise KeyError(f"no shape containing {substr!r}")


def set_para(tf, match, new_text):
    for para in tf.paragraphs:
        if match in para.text:
            if para.runs:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.text = new_text
            return


def add_para_after(tf, match, new_text):
    from pptx.text.text import _Paragraph
    for para in tf.paragraphs:
        if match in para.text:
            new_p = copy.deepcopy(para._p)
            para._p.addnext(new_p)
            np = _Paragraph(new_p, para._parent)
            if np.runs:
                np.runs[0].text = new_text
                for r in np.runs[1:]:
                    r.text = ""
            else:
                np.text = new_text
            return np


def add_note(s, text):
    tf = s.notes_slide.notes_text_frame
    ex = tf.text
    tf.text = (ex + "\n\n" if ex.strip() else "") + "[2026-07-12 rev3] " + text


def copy_shape(src_slide, dst_slide, substr):
    for sh in src_slide.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            el = copy.deepcopy(sh._element)
            dst_slide.shapes._spTree.append(el)
            for sh2 in dst_slide.shapes:
                if sh2._element is el:
                    return sh2


def slide_index(p, s):
    for i, x in enumerate(p.slides):
        if x == s:
            return i
    raise KeyError


def delete_slide(p, s):
    lst = p.slides._sldIdLst
    lst.remove(list(lst)[slide_index(p, s)])


def move_after(p, s, target):
    lst = p.slides._sldIdLst
    el = list(lst)[slide_index(p, s)]
    lst.remove(el)
    tpos = slide_index(p, target)
    lst.insert(tpos + 1, el)


def run(path):
    p = Presentation(path)

    # 129 / 12A: drop the per-state adaptation tables
    for anchor in ("5% review budget: national as-is vs adapted",
                   "10% review budget: national as-is vs adapted"):
        s = find_slide(p, anchor, required=False)
        if s is not None:
            delete_slide(p, s)

    # 135: engines teaching diagram, right after the intuition slide
    if find_slide(p, "How the two engines complement", required=False) is None:
        s3 = find_slide(p, "Intuition for new modeling approach")
        s_d = p.slides.add_slide(s3.slide_layout)
        t_d = copy_shape(s3, s_d, "Intuition for new modeling approach")
        set_para(t_d.text_frame, "Intuition for new modeling approach",
                 "How the two engines complement each other")
        s_d.shapes.add_picture(DIAGRAM, Inches(0.65), Inches(1.15),
                               width=Inches(8.7))
        add_note(s_d, "Comment 'find or create a diagram of a (uniform) random "
                 "forest... complementary to xgboost': schematic created "
                 "(methods/draw_engines_concept_diagram.R). Terminology note: "
                 "our measured engine is ranger with mtry=2 (constrained "
                 "splits), which captures the same teaching point as uniform "
                 "random forests - forced variety per split; uniform RFs "
                 "themselves are untested here.")
        move_after(p, s_d, s3)

    # 12B: clarify the few-dozen slide and place it after the blend intro
    s_few = find_slide(p, "The deployed list is a few dozen rules")
    tf_few = find_tf(s_few, "budget fill scans")
    if "blended lists behave the same" not in tf_few.text:
        add_para_after(tf_few, "budget fill scans",
                       "(Shown for the national pool; the blended lists behave "
                       "the same - median 23 / 42 rules run at 5% / 10%.)")
        add_note(s_few, "Comment 'is this blend? what is this?': it is the "
                 "budget-fill mechanics, measured on the national pool; the "
                 "blended lists inherit them (median activated 23 at 5%, 42 at "
                 "10% - blended_frozen_results.csv). Clarifying line added and "
                 "slide moved after the blend introduction.")

    # 130 + 13A + 12B: blend section order:
    # blend intro -> PR 10% (second slide) -> few dozen -> deliverable ->
    # freeze -> handoff
    s_intro = find_slide(p, "Blending state and national rules")
    s_pr = find_slide(p, "Blend vs national vs own rules, 10% budget")
    s_deliv = find_slide(p, "The deliverable: one ranked rule list")
    s_freeze = find_slide(p, "Freeze each state's list in advance")
    s_hand = find_slide(p, "What a state is handed")
    s_rules = find_slide(p, "The most widely deployed rules")
    move_after(p, s_intro, s_rules)
    move_after(p, s_pr, s_intro)
    move_after(p, s_few, s_pr)
    move_after(p, s_deliv, s_few)
    move_after(p, s_freeze, s_deliv)
    move_after(p, s_hand, s_freeze)
    add_note(s_intro, "Comment 'move this to the top of the blend section': "
             "blend section now runs intro -> per-state chart -> few-dozen "
             "mechanics -> deliverable -> freeze cost -> handoff table.")

    # 13A: 5% blend chart into the appendix
    if find_slide(p, "Appendix: blend vs national vs own", required=False) is None:
        s_app = find_slide(p, "Appendix: one-at-a-time xgboost sweep")
        s_x = p.slides.add_slide(s_app.slide_layout)
        t_x = copy_shape(s_app, s_x, "Appendix: one-at-a-time xgboost")
        set_para(t_x.text_frame, "Appendix: one-at-a-time xgboost",
                 "Appendix: blend vs national vs own, 5% budget")
        s_x.shapes.add_picture(
            os.path.join(DEPLOY, "blend_vs_state_vs_national_budget05.png"),
            Inches(2.9), Inches(0.75), width=Inches(4.3))
        add_note(s_x, "Comment: 5% version of the blend per-state chart, "
                 "placed in the appendix as requested.")

    p.save(path)
    print(f"updated {path} ({len(p.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    run(REV_A)
    run(REV_B)
