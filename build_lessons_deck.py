# Builds lessons_getting_more_signal_from_snap_qc_data.pptx from the how_to
# deck (theme + text-box formatting cloned from its appendix slides). Each
# lesson is standalone: the title is the advice, the body says what happened
# and what to do, pipeline jargon kept out so non-users of the pipeline can
# apply it.
import copy
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt

SRC = "how_to_use_the_snap_qc_inclusion_rules_scripts.pptx"
OUT = "lessons_getting_more_signal_from_snap_qc_data.pptx"
TEMPLATE_IDX = 16          # "Appendix: lucky rules, measured" - title/body/tag boxes
TAG = "Lessons from the SNAP QC error-mining project - July 2026"

shutil.copy(SRC, OUT)
p = Presentation(OUT)
tmpl = p.slides[TEMPLATE_IDX]

# grab template shape XML for title / body / tag text boxes
shapes = [sh for sh in tmpl.shapes if sh.has_text_frame]
tmpl_title = next(sh for sh in shapes if sh.text_frame.text.startswith("Appendix:"))
tmpl_body  = next(sh for sh in shapes if sh.text_frame.text.startswith("- "))
tmpl_tag   = next(sh for sh in shapes if "pipeline evidence" in sh.text_frame.text)


def set_line(tf, text):
    first = tf.paragraphs[0]
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    if first.runs:
        first.runs[0].text = text
        for r in first.runs[1:]:
            r.text = ""
    else:
        first.text = text


def set_body(tf, lines):
    set_line(tf, lines[0])
    first_p = tf.paragraphs[0]._p
    for line in lines[1:]:
        new_p = copy.deepcopy(first_p)
        first_p.getparent().addnext(new_p)
        first_p = new_p
    for para, line in zip(tf.paragraphs, lines):
        if para.runs:
            para.runs[0].text = line
            for r in para.runs[1:]:
                r.text = ""


def add_lesson(title, bullets, figure=None, fig_aspect=None, table=None):
    s = p.slides.add_slide(tmpl.slide_layout)
    for src_sh in (tmpl_title, tmpl_body, tmpl_tag):
        s._element.spTree.append(copy.deepcopy(src_sh._element))
    boxes = [sh for sh in s.shapes if sh.has_text_frame]
    title_sh = next(sh for sh in boxes if sh.text_frame.text.startswith("Appendix:"))
    body_sh  = next(sh for sh in boxes if sh.text_frame.text.startswith("- "))
    tag_sh   = next(sh for sh in boxes if "pipeline evidence" in sh.text_frame.text)
    set_line(title_sh.text_frame, title)
    set_body(body_sh.text_frame, bullets)
    set_line(tag_sh.text_frame, TAG)
    # estimate wrapped lines (body box is ~9.1in wide, ~105 chars per line)
    wrapped = sum(max(1, -(-len(b) // 105)) for b in bullets)
    fig_top = 0.85 + 0.28 * wrapped + 0.2
    if figure:
        aspect = fig_aspect if fig_aspect else 5 / 8
        max_h = 5.05 - fig_top
        h = min(max_h, 3.4)
        w = h / aspect
        if w > 9.2:
            w = 9.2
            h = w * aspect
        s.shapes.add_picture(figure, Inches((10 - w) / 2), Inches(fig_top),
                             width=Inches(w), height=Inches(h))
    if table:
        rows, cols = len(table), len(table[0])
        tw = Inches(6.5)
        th = Inches(0.32 * rows)
        gt = s.shapes.add_table(rows, cols, Inches((10 - 6.5) / 2),
                                Inches(fig_top), tw, th).table
        for i, row in enumerate(table):
            for j, val in enumerate(row):
                c = gt.cell(i, j)
                c.text = val
                for para in c.text_frame.paragraphs:
                    for r in para.runs:
                        r.font.size = Pt(12)
                        r.font.bold = (i == 0)
    return s


n_before = len(p.slides._sldIdLst)

# ── title slide ──────────────────────────────────────────────────────────────
s = p.slides.add_slide(tmpl.slide_layout)
for src_sh in (tmpl_title, tmpl_body):
    s._element.spTree.append(copy.deepcopy(src_sh._element))
boxes = [sh for sh in s.shapes if sh.has_text_frame]
t_sh = next(sh for sh in boxes if sh.text_frame.text.startswith("Appendix:"))
b_sh = next(sh for sh in boxes if sh.text_frame.text.startswith("- "))
set_line(t_sh.text_frame, "Getting more signal out of SNAP QC data")
set_body(b_sh.text_frame, [
    "Lessons from building and rebuilding an error-mining pipeline on the public QC files.",
    "Written to stand alone: each lesson applies to any analysis of these data, whether or not you use our pipeline.",
    "Code and evidence: https://github.com/giannella/snap_qc (modeling_findings.md)"])

# ── context ──────────────────────────────────────────────────────────────────
add_lesson(
    "Where these lessons come from",
    ["- We mine the public SNAP QC files for interpretable rules that flag high-risk cases for review (or safely exclude low-risk ones).",
     "- In one week we rebuilt the data pipeline, re-validated the statistics, and tested rule transfer across states.",
     "- Several 'modeling' problems turned out to be DATA problems, and several data problems were invisible until we measured them. The lessons below are ordered data-first."])

# ── data lessons ─────────────────────────────────────────────────────────────
add_lesson(
    "Know what your data cannot see",
    ["- The public QC files EXCLUDE cases found ineligible - which are 100%-of-benefit errors. The technical documentation quantifies the exclusions; almost nobody reads those counts.",
     "- Share of each state's error population visible in the public files (2022-24): national 71%. Some states are far worse.",
     "- Before promising results from any public-data model, compute this number for your state. Below ~60%, treat public-data results as a supplement and run your analysis on internal data, which contains the ineligible determinations."],
    table=[["state", "visible share of errors"],
           ["New Jersey", "43%"],
           ["Tennessee", "51%"],
           ["Arkansas / Missouri / Utah", "~53%"],
           ["Washington / Virginia / Louisiana", "78-81%"],
           ["national", "71%"]])

add_lesson(
    "Missing usually means 'not claimed', not 'unknown'",
    ["- Deduction fields (dependent care, medical, child support, homeless) are blank when the household didn't claim the deduction. Dropping rows with blanks silently deletes real, clean cases.",
     "- In Washington, dropping deduction blanks removed ~16% of the caseload before we caught it.",
     "- Zero-fill these fields and keep an 'imputed' flag. Reserve row-dropping for fields where blank truly means unknown (for us: rent and utilities)."])

add_lesson(
    "Don't drop cases you can't fully reconcile - flag them",
    ["- QC reviews can find MULTIPLE errors in one case. Cases with a second error element didn't fit our single-error reconstruction, so an early version dropped them: 31% of all error cases, silently.",
     "- Keeping them (with a flag) tripled the number of reliable patterns we could find - not because multi-error cases are special, but because a third more error data tightens every statistical bound.",
     "- The flag also let us verify the fix: new patterns were NOT concentrated on multi-error cases (34% of their catches vs 32% base) - the gain was statistical power, not a new error type."])

add_lesson(
    "Rebuild derived datasets from the script, never by hand",
    ["- Our modelling frame was once saved by hand from an interactive session. It silently descended from a version with the multi-error drop still active - and every result for weeks was mined on 69% of the true errors.",
     "- The fix is structural, not behavioral: the build script itself saves the frame as its last step, so the saved data can never drift from the code that claims to produce it.",
     "- When results move after a data rebuild, diff the OUTPUTS: we classify every mined rule as exact / threshold-shifted / overlapping / dropped / new, so a data change's fingerprint is auditable."])

# ── statistics lessons ───────────────────────────────────────────────────────
add_lesson(
    "Screening many candidates on raw performance selects lucky ones",
    ["- Generate 100,000 candidate rules and keep those with the best measured precision, and you mostly keep rules that got lucky: our 'train precision >= 0.20' shortlist delivered only ~0.10 on a year it had never seen.",
     "- The estimates were nearly unbiased BEFORE selection (r = 0.83 train vs holdout). Selection itself creates the bias - textbook regression to the mean, not overfitting.",
     "- This applies to any 'fit many, keep the best' workflow: feature screens, subgroup analyses, model leaderboards."],
    figure="presentation_figures/winners_curse_raw_vs_lcb.png", fig_aspect=1.0)

add_lesson(
    "Filter on a lower confidence bound, not the point estimate",
    ["- Keep a candidate only if the LOWER end of its precision confidence interval clears your bar. Small-sample lucky streaks fail this test automatically; well-supported patterns pass.",
     "- Same rule pool, three floor definitions: raw-precision floors overpromise (pick 0.50, get ~0.34) even after a confidence gate; lower-bound floors deliver at or above the promise (pick 0.30, get 0.38).",
     "- The lower bound is the only menu axis whose number means 'at least this'."],
    figure="presentation_figures/floor_definitions_educational.png", fig_aspect=9.5 / 13)

add_lesson(
    "Generate aggressively; let statistics do the vetoing",
    ["- Bigger ensembles (1000 trees instead of 100) do not trace a better precision-recall frontier - but they find several times more DISTINCT patterns at every quality bar.",
     "- That surplus is operational freedom: reviewers can veto patterns they distrust and still have substitutes.",
     "- The stringent lower-bound filter is what makes this safe: it removes the selection noise that large candidate pools would otherwise inject."],
    figure="presentation_figures/mine_big_filter_stringently.png", fig_aspect=5 / 8)

add_lesson(
    "Count every error you catch, not just the kind you looked for",
    ["- A rule built to find earned-income errors also flags cases carrying other error types - and those reviews succeed too.",
     "- Scored only against its intended error type, our earned-income rule set looked like 8% precision; scored against every error actually caught, 19%. The second number is what reviewers experience.",
     "- Whenever you model one outcome inside a family of related outcomes, report performance against the family too - the narrow metric can understate real performance by 2x."],
    figure="inclusion_rules_by_hh_size_v2/earned_income_lcb_sweep.png", fig_aspect=0.5)

add_lesson(
    "Look for the unmodeled majority",
    ["- Errors in deductions, shelter costs, and household composition ('other' errors) are the LARGEST category - 2,007 of 4,460 above-threshold errors in 2023 - and years of prior work never modeled them because they seemed heterogeneous.",
     "- They turned out to be learnable: this category now contributes the largest block of reliable patterns.",
     "- Inventory your outcome categories by SIZE before deciding what is modelable. 'Messy' is not the same as 'random'."],
    figure="inclusion_rules_by_hh_size_v2/other_error_lcb_sweep.png", fig_aspect=0.5)

add_lesson(
    "Split by structure - but split coarsely",
    ["- Household size changes which variables mean what (income per member, deduction scales), so we model size groups 1 / 2-3 / 4+ separately. Coarse splits kept enough data per group; a 5-way split performed worse.",
     "- Not every important group needs its own model: elderly/disabled households have a very different ERROR MIX, but an indicator variable let the models carve the caseload themselves - a separate stratum bought nothing.",
     "- Test 'feature vs stratum' empirically; the answer is usually 'feature' unless the group changes the meaning of other variables."],
    figure="presentation_figures/esap_error_mix.png", fig_aspect=5 / 8)

# ── small-sample and transfer lessons ────────────────────────────────────────
add_lesson(
    "Small samples need hard support floors, not just confidence bounds",
    ["- At national scale (100k+ cases), a stringent lower-bound filter alone controls selection noise. At single-state scale it did NOT: rules selected with the bound but tiny support had median holdout precision of ZERO.",
     "- Adding a hard floor - a rule must fire on at least ~30 training cases - changed failure from collapse to gentle deflation (~1/3).",
     "- Confidence bounds assume the model class is honest; tiny-support rules are where that assumption breaks first."])

add_lesson(
    "Match eras: recent similar data beats more mixed-era data",
    ["- For a state with weak local signal (Louisiana), pooling its own data across 2017-19 + 2022-24 made things WORSE (median holdout precision fell to zero) - error-generating processes drift with policy changes.",
     "- Training on five SIMILAR states' 2022-24 data - never touching Louisiana - delivered usable rules there (14% precision at 49% of error dollars, in a state where local mining had collapsed).",
     "- When data is thin, the instinct is to add years. Add NEIGHBORS from the same era instead."])

add_lesson(
    "The microdata already documents each state's policy choices",
    ["- Key policy options - broad-based categorical eligibility, reporting systems, certification periods, standard medical deduction, standard utility allowances, SSI CAP - are all readable off the QC microdata itself, per state and era.",
     "- That means state similarity can be MEASURED, not asserted: we combine policy vectors with 'which risk patterns fire here' profiles, weighting rare shared patterns more heavily (inverse-frequency).",
     "- Use similarity to pick donor states for thin-data problems - and recompute it per era, because states change policies."])

add_lesson(
    "Takeaways",
    ["- Data first: measure what your data cannot see, read missingness semantics, keep hard cases with flags, and make scripts - not hands - produce datasets.",
     "- Statistics second: any 'generate many, keep the best' workflow needs lower-bound filtering, honest holdout years, and support floors at small scale.",
     "- Scope last: score against all the value you create, hunt the unmodeled majority, and borrow strength from measured similarity when your own data runs out.",
     "- Everything here is reproducible from the public files: github.com/giannella/snap_qc"])

# ── drop the original how_to slides ──────────────────────────────────────────
sldIdLst = p.slides._sldIdLst
for sld in list(sldIdLst)[:n_before]:
    sldIdLst.remove(sld)

p.save(OUT)
print("built", OUT, "with", len(p.slides._sldIdLst), "slides")
